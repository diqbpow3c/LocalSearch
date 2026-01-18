# This is an app that performs hybrid search (BM25 + embedding similarity) on local files and folders.
# It uses FAISS for embedding similarity and rank_bm25 for BM25 search.
# The app has a GUI built with PySide6.


import os, sys
os.environ["DO_NOT_TRACK"] = "true"  # for speeding up unstructured
os.environ["HF_HUB_OFFLINE"] = "1"
import logging, json
from datetime import datetime, timedelta
from PySide6.QtWidgets import (QApplication, QMenu, QMainWindow,
                               QWidget, QVBoxLayout, QHBoxLayout, QPushButton,
                               QLineEdit, QListWidget, QTableWidget, QSystemTrayIcon,
                               QTableWidgetItem, QHeaderView, QComboBox, QFileDialog,
                               QTabWidget, QSpinBox, QStatusBar, QMessageBox,
                               QLabel, QAbstractItemView, QTextEdit, QCheckBox,
                               QSizePolicy, QScrollArea, QSplitter,
                               QStackedWidget, QTextBrowser, )
from PySide6.QtCore import (Slot, Qt, QUrl, QSettings, QByteArray, QPropertyAnimation, QRect)
from PySide6.QtGui import (Qt, QAction, QFont, QFontDatabase, QDesktopServices,
                           QIcon, QPixmap, QColor, QImage, QPainter)

from pathlib import Path
import configs
from configs import *
from utils import highlight_keywords_in_text

os.makedirs(INDEX_DIR, exist_ok=True)
# Workaround that makes pyinstaller exe work normally
log_file = open(LOG_PATH, "w", encoding="utf-8")
sys.stdout = log_file
sys.stderr = log_file
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        # logging.FileHandler(LOG_PATH, encoding="utf-8"),
        logging.StreamHandler(sys.stdout),
    ],
)


logger = logging.getLogger(__name__)

from misc_threads import *


class UniversalPreviewPane(QStackedWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        # try to hide the importing latency behind the splash screen
        from pymupdf import open as pymupdf_open  # PyMuPDF
        from markdown import markdown as mdmd
        from docx import Document
        from pptx import Presentation

        self.Document = Document
        self.Presentation = Presentation
        self.pymupdf_open = pymupdf_open
        self.mdmd = mdmd

        # 1. Image & PDF Viewer (via Pixmap)
        self.image_viewer = QLabel()
        self.image_viewer.setAlignment(Qt.AlignCenter)
        self.image_scroll = QScrollArea()
        self.image_scroll.setWidgetResizable(True)
        self.image_scroll.setWidget(self.image_viewer)

        # 2. Markdown & Rich Text Viewer
        self.text_browser = QTextBrowser()
        self.text_browser.setOpenExternalLinks(True)

        # 3. Plain Text / Code Viewer
        self.code_viewer = QTextEdit()
        self.code_viewer.setReadOnly(True)
        self.code_viewer.setLineWrapMode(QTextEdit.NoWrap)

        # 4. Error/Fallback View
        self.error_label = QLabel("Preview not available for this file type.")
        self.error_label.setAlignment(Qt.AlignCenter)

        # Add all to the stack
        self.addWidget(self.image_scroll)  # Index 0
        self.addWidget(self.text_browser)  # Index 1
        self.addWidget(self.code_viewer)  # Index 2
        self.addWidget(self.error_label)  # Index 3

    def handle_preview(self, file_path):
        """Main entry point to update the preview based on file path."""
        if not os.path.exists(file_path):
            self.show_text("File not found.")
            return
        MAX_FILE_PREVIEW_SIZE_MB = 20
        if os.path.getsize(file_path) > MAX_FILE_PREVIEW_SIZE_MB * 1024 * 1024:  # multi MB limit
            self.show_text(f"File too large (>{MAX_FILE_PREVIEW_SIZE_MB} MB) to preview.")
            return
        ext = os.path.splitext(file_path)[1].lower()

        try:
            if ext in [".jpg", ".jpeg", ".png", ".bmp", ".gif"]:
                self.preview_image(file_path)
            elif ext == ".pdf":
                self.preview_pdf(file_path)
            elif ext == ".md":
                self.preview_markdown(file_path)
            elif ext == ".docx":
                self.preview_docx(file_path)
            elif ext == ".pptx":
                self.preview_pptx(file_path)
            elif ext in [".txt", ".py", ".json", ".csv", ".log", ".xml"]:
                self.preview_text(file_path)
            else:
                self.show_text(f"No preview available for {ext}")
        except Exception as e:
            self.show_text(f"Error loading preview: {str(e)}")

    def preview_image(self, path):
        pixmap = QPixmap(path)
        self._set_pixmap_scaled(pixmap)
        self.setCurrentIndex(0)

    def preview_pdf(self, path):
        """Renders the first page of a PDF as an image."""
        doc = self.pymupdf_open(path)
        page = doc.load_page(0)
        pix = page.get_pixmap()
        fmt = QImage.Format_RGBA8888 if pix.alpha else QImage.Format_RGB888
        qimg = QImage(pix.samples, pix.width, pix.height, pix.stride, fmt)
        self._set_pixmap_scaled(QPixmap.fromImage(qimg))
        doc.close()
        self.setCurrentIndex(0)

    def preview_markdown(self, path):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            html = self.mdmd(f.read())
            # Add basic CSS to make it look nice
            styled_html = f"<style>body {{ font-family: sans-serif; padding: 10px; }}</style>{html}"
            self.text_browser.setHtml(styled_html)
        self.setCurrentIndex(1)

    def preview_docx(self, path):
        doc = self.Document(path)
        text = [p.text for p in doc.paragraphs[:50]]  # Limit to first 50 paragraphs
        self.text_browser.setPlainText("\n".join(text))
        self.setCurrentIndex(1)

    def preview_pptx(self, path):
        prs = self.Presentation(path)
        text = []
        for i, slide in enumerate(prs.slides[:5]):  # Limit to first 5 slides
            text.append(f"--- Slide {i + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        self.text_browser.setPlainText("\n".join(text))
        self.setCurrentIndex(1)

    def preview_text(self, path):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            self.code_viewer.setPlainText(f.read(10000))  # Load first 10k chars
        self.setCurrentIndex(2)

    def show_text(self, message):
        self.error_label.setText(message)
        self.setCurrentIndex(3)

    def clear(self):
        self.show_text("")

    def _set_pixmap_scaled(self, pixmap):
        """Helper to scale images to fit the scroll area width."""
        if not pixmap.isNull():
            # Adjust scaling logic based on your preference
            scaled = pixmap.scaled(
                self.size(), Qt.KeepAspectRatio, Qt.SmoothTransformation
            )
            self.image_viewer.setPixmap(scaled)


class SettingsTab(QWidget):
    def __init__(self, sercher_instance):
        super().__init__()
        self.settings = QApplication.instance().settings
        self.autostart_checkbox = QCheckBox("Start the application on system startup")
        self.autostart_checkbox.stateChanged.connect(self.on_autostart_checkbox_changed)
        self.help_button = QPushButton("Help")
        self.help_button.clicked.connect(self.on_help_button_pushed)
        self.help_window = None
        self.tray_checkbox = QCheckBox("Minimize to system tray on close")
        self.tray_checkbox.setChecked(
            self.settings.value("minimize_to_tray", True, type=bool)
        )
        self.tray_checkbox.stateChanged.connect(self.toggle_tray_setting)
        self.spin_box = QSpinBox()
        self.spin_box.setRange(8, 20)
        self.spin_box.setFixedWidth(80)
        self.searcher = sercher_instance
        self.init_ui()
        self.load_settings()

    def init_ui(self):
        layout = QVBoxLayout()
        layout.addWidget(self.autostart_checkbox)
        layout.addWidget(self.tray_checkbox)
        h_layout = QHBoxLayout()
        input_label = QLabel("Font size (takes effect after restart):")
        h_layout.addWidget(input_label)
        h_layout.addWidget(self.spin_box, alignment=Qt.AlignLeft)  # type: ignore
        h_layout.addStretch()
        layout.addLayout(h_layout)

        index_management_layout = QHBoxLayout()
        # Folder Paths
        folder_layout = QVBoxLayout()
        folder_layout.addWidget(QLabel("Indexed Folders"))
        self.folder_list = QListWidget()
        for path in self.searcher.folder_paths:
            self.folder_list.addItem(path)
        folder_layout.addWidget(self.folder_list)

        folder_buttons_layout = QHBoxLayout()
        self.add_folder_button = QPushButton("Add Folder")
        folder_buttons_layout.addWidget(self.add_folder_button)

        self.remove_folder_button = QPushButton("Remove Selected Folder")
        folder_buttons_layout.addWidget(self.remove_folder_button)
        folder_layout.addLayout(folder_buttons_layout)
        index_management_layout.addLayout(folder_layout)

        # File Paths
        file_layout = QVBoxLayout()
        file_layout.addWidget(QLabel("Indexed Files"))
        self.file_list = QListWidget()
        for path in self.searcher.file_paths:
            self.file_list.addItem(path)
        file_layout.addWidget(self.file_list)
        file_buttons_layout = QHBoxLayout()
        self.add_file_button = QPushButton("Add File")
        file_buttons_layout.addWidget(self.add_file_button)
        self.remove_file_button = QPushButton("Remove Selected File")
        file_buttons_layout.addWidget(self.remove_file_button)
        file_layout.addLayout(file_buttons_layout)
        index_management_layout.addLayout(file_layout)

        layout.addLayout(index_management_layout)

        self.rebuild_index_button = QPushButton("Rebuild Index Now")
        layout.addWidget(self.rebuild_index_button)
        layout.addWidget(self.help_button)
        misc_info_label = QLabel(
            f"Supported Extensions:\nFile content and file paths for the following extensions are indexed\n{CONTENT_PARSING_EXTENSIONS}\nOnly the file paths are indexed for the following extensions:\n{NO_CONTENT_PARSING_EXTENSIONS}\nNo OCR is used (e.g., for PDF).")
        layout.addWidget(misc_info_label)
        layout.addStretch(1)  # Pushes everything to the top
        self.setLayout(layout)

    def load_settings(self):
        autostart_enabled = self.settings.value("autostart_enabled", False, type=bool)
        self.autostart_checkbox.setChecked(autostart_enabled)
        _font_size = self.settings.value("font_size", 10, type=int)
        self.spin_box.setValue(_font_size)

    def toggle_tray_setting(self, state):
        minimize_to_tray = state == 2
        self.settings.setValue("minimize_to_tray", minimize_to_tray)

    def on_help_button_pushed(self):
        to_remove = [".", "[", "]", "'"]
        supported_types = str(SUPPORTED_EXTENSIONS)
        tmp = supported_types.translate(str.maketrans("", "", "".join(to_remove)))
        help_text = f"""Usage：
                    1.First, you need to add folders or single files to be indexed in the settings tab
                    2.Double clicking the file name in the search results would open the file, and double clicking the folder path would open the folder
                    """

        help_text = "\n".join(
            line.lstrip() for line in help_text.splitlines()
        )  # remove indentation
        QMessageBox.information(self, "Help", help_text)

    def on_autostart_checkbox_changed(self, state):
        is_checked = state == Qt.CheckState.Checked.value
        self.settings.setValue("autostart_enabled", self.autostart_checkbox.isChecked())  # Save the new state immediately

        if is_checked:
            logger.info("Autostart checkbox checked. Attempting to enable autostart...")
            self.enable_autostart()
        else:
            logger.info(
                "Autostart checkbox unchecked. Attempting to disable autostart..."
            )
            self.disable_autostart()

    def enable_autostart(self):
        script_path = Path(__file__).resolve()
        app_path = f'"{sys.executable}"'  # Or path to your .exe if frozen
        if is_running_in_pyinstaller():
            command = app_path + " --autostart"
        else:
            command = f'"{sys.executable}" "{script_path}" --autostart'
        if sys.platform.startswith("win"):
            import winreg

            try:
                key = winreg.OpenKey(
                    winreg.HKEY_CURRENT_USER,
                    r"Software\Microsoft\Windows\CurrentVersion\Run",
                    0,
                    winreg.KEY_SET_VALUE,
                )
                # Get the path to your executable. This might need to be more robust.
                # For a frozen executable, sys.executable points to the executable.
                # For a script, you might need to run python and the script.

                winreg.SetValueEx(
                    key, "LocalSearch", 0, winreg.REG_SZ, command
                )
                winreg.CloseKey(key)
                logger.info("Autostart enabled in Windows Registry.")
            except Exception as e:
                QMessageBox.warning(
                    self,
                    "Autostart Error",
                    f"Failed to enable autostart on Windows: {e}",
                )
                logger.info(f"Windows autostart error: {e}")

        elif sys.platform.startswith("linux"):
            # Linux: Create a .desktop file in ~/.config/autostart/
            autostart_dir = os.path.expanduser("~/.config/autostart")
            os.makedirs(autostart_dir, exist_ok=True)
            desktop_file_path = os.path.join(autostart_dir, "LocalSearch.desktop")
            # Get the path to your executable.
            # For a frozen app, sys.executable is the app path.
            # For a script, you'd need the path to the script and python executable.
            exec_path = (
                sys.executable
            )  # Or the path to your main script, e.g., f"python3 {os.path.abspath(__file__)}"

            desktop_content = f"""[Desktop Entry]
            Type=Application
            Exec={command}
            Hidden=false
            NoDisplay=false
            X-GNOME-Autostart-enabled=true
            Name=MyProgram
            Comment=Start MyProgram on boot
            """
            try:
                with open(desktop_file_path, "w") as f:
                    f.write(desktop_content)
                logger.info(f"Autostart .desktop file created at: {desktop_file_path}")
            except Exception as e:
                QMessageBox.warning(
                    self, "Autostart Error", f"Failed to enable autostart on Linux: {e}"
                )
                logger.info(f"Linux autostart error: {e}")
        else:
            QMessageBox.warning(
                self,
                "Autostart",
                f"Autostart not supported on this operating system ({sys.platform}).",
            )

    def disable_autostart(self):

        if sys.platform.startswith("win"):
            # Windows: Remove the registry entry or delete the shortcut.
            # QMessageBox.information(self, "Autostart",
            #                         "On Windows, you'd typically remove the registry entry "
            #                         "or delete the shortcut from the Startup folder.")
            logger.info("Windows disable autostart.")
            import winreg

            try:
                key = winreg.OpenKey(
                    winreg.HKEY_CURRENT_USER,
                    r"Software\Microsoft\Windows\CurrentVersion\Run",
                    0,
                    winreg.KEY_SET_VALUE,
                )
                winreg.DeleteValue(key, "LocalSearch")
                winreg.CloseKey(key)
                logger.info("Autostart disabled in Windows Registry.")
            except FileNotFoundError:
                logger.info(
                    "Autostart entry not found in Windows Registry (already disabled?)."
                )
            except Exception as e:
                QMessageBox.warning(
                    self,
                    "Autostart Error",
                    f"Failed to disable autostart on Windows: {e}",
                )
                logger.info(f"Windows disable autostart error: {e}")

        elif sys.platform.startswith("linux"):
            # Linux: Delete the .desktop file.
            desktop_file_path = os.path.expanduser(
                "~/.config/autostart/LocalSearch.desktop"
            )
            try:
                if os.path.exists(desktop_file_path):
                    os.remove(desktop_file_path)
                    logger.info(
                        f"Autostart .desktop file removed from: {desktop_file_path}"
                    )
                else:
                    logger.info(
                        "Autostart .desktop file not found (already disabled?)."
                    )
            except Exception as e:
                QMessageBox.warning(
                    self,
                    "Autostart Error",
                    f"Failed to disable autostart on Linux: {e}",
                )
                logger.info(f"Linux disable autostart error: {e}")
        else:
            pass  # Already handled by enable_autostart message


class MainWindow(QMainWindow):
    """
    The main application window for the GUI.
    """

    def __init__(self):
        super().__init__()
        from searcher_onnx import Searcher  # move the importing of heavy libraries here for faster program initialization
        self.settings = QApplication.instance().settings
        self.preview_pane = UniversalPreviewPane()
        self.setWindowTitle("LocalSearch")
        geometry = self.settings.value("window_geometry", QByteArray())
        if not geometry.isEmpty():
            self.restoreGeometry(geometry)
        else:
            self.setGeometry(100, 100, 800, 600)  # Initial window size
        # Ensure the main window is resizable: set a reasonable minimum and no restrictive maximum
        self.setMinimumSize(600, 400)
        self.setMaximumSize(16777215, 16777215)
        # Allow central widget to expand so resizing the window changes height/width
        self.central_widget = QWidget()
        self.setCentralWidget(self.central_widget)
        self.central_widget.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        # recreate main layout after resetting central widget
        self.main_layout = QVBoxLayout(self.central_widget)
        icon = QIcon("./resources/icon.png")
        self.setWindowIcon(icon)

        self.tray_icon = QSystemTrayIcon(icon, self)
        self.tray_icon.setToolTip("LocalSearch")
        tray_menu = QMenu()
        restore_action = QAction("Restore Window", self)
        quit_action = QAction("Exit", self)
        restore_action.triggered.connect(self.show)
        quit_action.triggered.connect(self.close_app)
        tray_menu.addAction(restore_action)
        tray_menu.addAction(quit_action)
        self.tray_icon.setContextMenu(tray_menu)
        self.tray_icon.activated.connect(self.on_tray_icon_activated)
        self.tray_icon.show()

        if configs.DEVICE=='gpu':
            try:
                self.searcher = Searcher()
            except Exception as e:
                logger.error(e)
                QMessageBox.warning(self, "GPU initialization error", f"Error {e}. You need to install torch (GPU version) along with onnxruntime-gpu correctly.")
        else:
            self.searcher = Searcher()

        self.searcher.index_rebuild_progress_signal.connect(
            self.on_rebuild_signal_progress
        )
        self.index_rebuilder_thread = None
        self.file_monitor_thread = None
        # Store last search results so filters can be applied locally without re-searching
        self.last_search_results = None
        self.last_highlight_tokens = None
        self.last_query = ""

        self._load_config()  # Load paths and last search mode
        self.searcher.file_paths = set(self.config.get("file_paths", []))
        self.searcher.folder_paths = set(self.config.get("folder_paths", []))

        self._setup_ui()
        self._setup_threads()  # Initialize threads but don't start file_monitor_thread yet

        # Attempt to load index on startup
        if self.searcher.load_index_components():
            self.status_bar.showMessage("Index loaded successfully.")
            # Start file monitor after successful load
            self.file_monitor_thread.start()
        else:
            # If load fails or no index exists, rebuild
            self._start_rebuild_index()

        # show the help window on first launch
        is_first_time_launch = self.settings.value("is_first_time_launch", 1)
        if is_first_time_launch:
            self.settings_tab.on_help_button_pushed()
            self.settings.setValue("is_first_time_launch", 0)

    def _load_config(self):
        """Loads configuration from a JSON file."""
        self.config = {}
        if os.path.exists(CONFIG_FILE):
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                self.config = json.load(f)
        else:
            logger.info(f"Config file not found: {CONFIG_FILE}")

    def _save_config(self):
        """Saves current configuration to a JSON file."""
        self.config["file_paths"] = list(self.searcher.file_paths)
        self.config["folder_paths"] = list(self.searcher.folder_paths)
        last_search_mode = list(SEARCH_MODE_MAPPING.values())[
            self.search_mode_combo.currentIndex()
        ]
        self.config["last_search_mode"] = last_search_mode
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump(self.config, f, indent=4, ensure_ascii=False)

    def _setup_ui(self):
        """Sets up the main user interface."""
        self.tab_widget = QTabWidget()
        self.main_layout.addWidget(self.tab_widget)

        self._setup_search_tab()
        self._setup_settings_tab()
        self.status_bar = QStatusBar()
        self.app_status_info = QLabel(f"Version {APP_VERSION}, Device: {configs.DEVICE}")
        self.app_status_info.setAlignment(Qt.AlignRight | Qt.AlignVCenter)
        self.status_bar.addPermanentWidget(self.app_status_info)
        self.setStatusBar(self.status_bar)
        self.status_bar.showMessage("Ready")


    def on_dedup_checkbox_toggled(self, state):
        is_checked = state == Qt.CheckState.Checked.value
        self.settings.setValue("dedup_checkbox_checked",
                               self.dedup_checkbox.isChecked())  # Save the new state immediately


    def _setup_search_tab(self):
        """Sets up the Search tab UI."""
        self.search_tab = QWidget()
        self.tab_widget.addTab(self.search_tab, "Search")
        search_layout = QVBoxLayout(self.search_tab)

        # Search Input and Mode
        search_input_layout = QHBoxLayout()
        self.query_input = QLineEdit()
        self.query_input.setClearButtonEnabled(True)
        self.query_input.setPlaceholderText("Enter your query...")
        self.query_input.returnPressed.connect(self._perform_search)
        search_input_layout.addWidget(self.query_input)

        self.search_mode_combo = QComboBox()
        self.search_mode_combo.addItems(list(SEARCH_MODE_MAPPING.keys()))
        # Set last used search mode from config
        last_mode = self.config.get("last_search_mode", "hybrid")
        index = list(SEARCH_MODE_MAPPING.values()).index(last_mode)
        self.search_mode_combo.setCurrentIndex(index)
        search_input_layout.addWidget(self.search_mode_combo)

        # checkbox for whether to perform deduplication of chunks for the same file in the results table
        self.dedup_checkbox = QCheckBox("Deduplicate")
        self.dedup_checkbox.setToolTip(
            "<span>When enabled, only the text chunk with the highest score will be shown for each document, instead of potentially showing many entries for one document. Takes effect for the next search.</span>"
        )
        dedup = self.settings.value("dedup_checkbox_checked", False, type=bool)
        self.dedup_checkbox.setChecked(dedup)
        self.dedup_checkbox.toggled.connect(self.on_dedup_checkbox_toggled)
        search_input_layout.addWidget(self.dedup_checkbox)

        self.search_button = QPushButton("Search")
        self.search_button.clicked.connect(self._perform_search)
        search_input_layout.addWidget(self.search_button)

        search_layout.addLayout(search_input_layout)

        # Main content area with sidebar and results table
        content_layout = QHBoxLayout()

        # Store default sidebar width for animation
        self.sidebar_default_width = 200

        # Create toggle button early so it can be used in results area
        self.toggle_sidebar_button = QPushButton("◀")  # Left arrow to collapse
        self.toggle_sidebar_button.setMaximumWidth(30)
        self.toggle_sidebar_button.clicked.connect(self._toggle_sidebar)

        # Left Sidebar for Filters (wrapped in a scroll area so items can be scrolled when space is small)
        self.sidebar = QScrollArea()
        self.sidebar.setWidgetResizable(True)
        sidebar_content = QWidget()
        sidebar_layout = QVBoxLayout(sidebar_content)
        # make sidebar more compact
        sidebar_layout.setSpacing(1)
        sidebar_layout.setContentsMargins(1, 1, 1, 1)

        # Sidebar header with toggle button
        sidebar_header_layout = QHBoxLayout()
        sidebar_title = QLabel("Filters")
        sidebar_title.setStyleSheet("font-weight: bold; font-size: 12px;")
        sidebar_header_layout.addWidget(sidebar_title)
        sidebar_header_layout.addStretch()
        sidebar_layout.addLayout(sidebar_header_layout)
        sidebar_layout.addSpacing(5)

        # Date Filter Section
        date_label = QLabel("Filter by Date:")
        date_label.setStyleSheet("font-weight: bold;")
        sidebar_layout.addWidget(date_label)

        self.date_filters = {}
        date_options = [
            ("All Time", None),
            ("Past 24 Hours", timedelta(hours=24)),
            ("Past 7 Days", timedelta(days=7)),
            ("Past 30 Days", timedelta(days=30)),
            ("Past 90 Days", timedelta(days=90)),
        ]

        for label, delta in date_options:
            checkbox = QCheckBox(label)
            checkbox.stateChanged.connect(self._on_filter_changed)
            self.date_filters[label] = (checkbox, delta)
            sidebar_layout.addWidget(checkbox)

        sidebar_layout.addSpacing(1)

        # File Type Filter Section
        file_type_label = QLabel("Filter by File Type:")
        file_type_label.setStyleSheet("font-weight: bold;")
        sidebar_layout.addWidget(file_type_label)

        self.file_type_filters = {}
        file_types = [
            ("All Types", None),
            ("PowerPoint", [".pptx"]),
            ("Excel", [".xlsx", ".xls"]),
            ("Word", [".docx"]),
            ("PDF", [".pdf"]),
            ("Text", [".txt"]),
            ("Markdown", [".md"]),
            ("CSV", [".csv"]),
            ("HTML", [".html", ".htm"]),
            ("ODT", [".odt"]),
            ("XML", [".xml"]),
        ]

        for label, extensions in file_types:
            checkbox = QCheckBox(label)
            checkbox.stateChanged.connect(self._on_filter_changed)
            self.file_type_filters[label] = (checkbox, extensions)
            sidebar_layout.addWidget(checkbox)

        # Set "All Time" as default
        self.date_filters["All Time"][0].setChecked(True)

        # Set "All Types" as default
        self.file_type_filters["All Types"][0].setChecked(True)

        sidebar_layout.addSpacing(1)

        # Clear Filters Button
        self.clear_filters_button = QPushButton("Clear Filters")
        self.clear_filters_button.clicked.connect(self._clear_filters)
        sidebar_layout.addWidget(self.clear_filters_button)

        sidebar_layout.addStretch()
        sidebar_content.setLayout(sidebar_layout)
        # Set sidebar width on the scroll area
        self.sidebar.setWidget(sidebar_content)
        self.sidebar.setMaximumWidth(240)

        # Restore sidebar visibility state from settings
        sidebar_visible = self.settings.value("sidebar_visible", True, type=bool)
        if not sidebar_visible:
            self.sidebar.setMaximumWidth(0)
            self.toggle_sidebar_button.setText("▶")

        # Create a container for sidebar with overlapping toggle button
        sidebar_container = QWidget()
        sidebar_container_layout = QVBoxLayout(sidebar_container)
        sidebar_container_layout.setContentsMargins(0, 0, 0, 0)
        sidebar_container_layout.setSpacing(0)
        sidebar_container_layout.addWidget(self.sidebar)

        # Add toggle button with negative margin to overlap sidebar edge
        sidebar_container_layout.addWidget(self.toggle_sidebar_button)

        content_layout.addWidget(sidebar_container)

        # Results area
        # results_area_layout = QHBoxLayout()

        # Search Results Table
        self.results_table = QTableWidget()
        table_headers = [
            "Filename",
            "Folder Path",
            "Snippet",
            "File Size",
            "Last Modified",
        ]
        self.results_table.setColumnCount(len(table_headers))
        self.results_table.setHorizontalHeaderLabels(table_headers)
        self.results_table.horizontalHeader().setSectionResizeMode(
            QHeaderView.ResizeMode.Interactive
        )
        self.results_table.horizontalHeader().setSectionsMovable(True)
        self.results_table.setWordWrap(True)
        header_state = self.settings.value("headerState")
        if header_state:
            self.results_table.horizontalHeader().restoreState(header_state)
        else:
            self.results_table.setColumnWidth(0, 180)
            self.results_table.setColumnWidth(1, 180)
            self.results_table.setColumnWidth(2, 280)

        # Resize columns to fit contents
        # self.results_table.resizeColumnsToContents()
        # self.results_table.setSortingEnabled(True) # Causes incomplete display
        self.results_table.setSortingEnabled(False)
        self.results_table.setSelectionBehavior(
            QAbstractItemView.SelectionBehavior.SelectItems
        )
        # self.results_table.setEditTriggers(QAbstractItemView.EditTrigger.NoEditTriggers)  # Make cells read-only
        self.results_table.setEditTriggers(
            QAbstractItemView.EditTrigger.SelectedClicked
        )
        self.results_table.cellDoubleClicked.connect(self._handle_table_double_click)
        # Enable custom context menu for result rows
        self.results_table.setContextMenuPolicy(Qt.CustomContextMenu)
        self.results_table.customContextMenuRequested.connect(
            self._show_results_context_menu
        )
        # Connect selection to preview
        self.results_table.itemSelectionChanged.connect(
            self._on_result_selection_changed
        )

        # Create a splitter for results table and preview pane (left=results, right=preview)
        self.splitter = QSplitter(Qt.Horizontal)
        self.splitter.addWidget(self.results_table)

        # Preview pane
        preview_container = QWidget()
        preview_layout = QVBoxLayout(preview_container)
        preview_label = QLabel("Preview")
        preview_label.setStyleSheet("font-weight: bold; font-size: 11px; margin: 5px;")
        preview_layout.addWidget(preview_label)

        self.preview_pane.setMinimumWidth(20)
        preview_layout.addWidget(self.preview_pane)
        preview_layout.setContentsMargins(0, 0, 0, 0)
        preview_layout.setSpacing(0)

        self.splitter.addWidget(preview_container)
        self.splitter.setStretchFactor(0, 1)
        self.splitter.setStretchFactor(1, 1)

        # Store splitter state
        splitter_state = self.settings.value("splitterState")
        if splitter_state:
            self.splitter.restoreState(splitter_state)

        content_layout.addWidget(self.splitter)
        content_layout.setStretch(1, 1)
        search_layout.addLayout(content_layout)

    def _setup_settings_tab(self):
        self.settings_tab = SettingsTab(self.searcher)
        self.tab_widget.addTab(self.settings_tab, "Settings")
        self.settings_tab.add_folder_button.clicked.connect(self._add_folder)
        self.settings_tab.remove_folder_button.clicked.connect(self._remove_folder)
        self.settings_tab.add_file_button.clicked.connect(self._add_file)
        self.settings_tab.remove_file_button.clicked.connect(self._remove_file)
        self.settings_tab.rebuild_index_button.clicked.connect(self._start_rebuild_index)

    def _setup_threads(self):
        """Sets up the background threads for indexing and file monitoring."""
        # Index Rebuilder Thread
        self.index_rebuilder_thread = IndexRebuilderThread(self.searcher)
        self.index_rebuilder_thread.rebuild_started.connect(self._on_rebuild_started)
        self.index_rebuilder_thread.rebuild_finished.connect(self._on_rebuild_finished)

        # File Monitor Thread (will be started after initial index load/rebuild)
        interval = self.settings.value("monitor_interval", FILE_CHANGE_CHECK_INTERVAL, type=int)
        self.file_monitor_thread = FileChangeMonitorThread(
            self.searcher, interval=interval
        )
        self.file_monitor_thread.changes_detected.connect(
            self._on_file_changes_detected
        )
        self.file_monitor_thread.monitoring_status_update.connect(
            self.status_bar.showMessage
        )

    @Slot()
    def _handle_table_double_click(self, row, column):
        item_widget = self.results_table.item(row, 1)
        item = item_widget.text() if item_widget is not None else None
        if not item:
            return
        if column == 1:  # open the folder
            if os.path.exists(item):
                QDesktopServices.openUrl(QUrl.fromLocalFile(item))
        else:  # open the file
            if os.path.exists(item):
                QDesktopServices.openUrl(QUrl.fromLocalFile(item))

    @Slot("QPoint")
    def _show_results_context_menu(self, pos):
        """Show context menu for a result row with actions: Open folder, Copy path."""
        row = self.results_table.rowAt(pos.y())
        if row < 0:
            return

        item_widget = self.results_table.item(row, 1)
        path = item_widget.text() if item_widget is not None else None
        if not path:
            return

        menu = QMenu(self)
        open_folder_act = QAction("Open folder", self)
        copy_path_act = QAction("Copy path", self)
        open_folder_act.triggered.connect(lambda: self._open_folder_for_row(path))
        copy_path_act.triggered.connect(lambda: self._copy_path_for_row(path))
        menu.addAction(open_folder_act)
        menu.addAction(copy_path_act)
        menu.exec(self.results_table.viewport().mapToGlobal(pos))

    def _open_folder_for_row(self, path):
        folder = str(Path(path))
        if os.path.exists(folder):
            QDesktopServices.openUrl(QUrl.fromLocalFile(folder))
        else:
            QMessageBox.warning(self, "Open folder", f"Path not found: {folder}")

    def _copy_path_for_row(self, path):
        try:
            QApplication.instance().clipboard().setText(path)
            self.status_bar.showMessage("Path copied to clipboard", 2000)
        except Exception:
            QMessageBox.information(self, "Copy path", path)

    @Slot()
    def _toggle_sidebar(self):
        """Toggle sidebar visibility with animation."""
        target_width = 0 if self.sidebar.width() > 50 else self.sidebar_default_width

        # Create animation for sidebar width
        self.sidebar_animation = QPropertyAnimation(self.sidebar, b"maximumWidth")
        self.sidebar_animation.setDuration(300)  # 300ms animation
        self.sidebar_animation.setStartValue(self.sidebar.maximumWidth())
        self.sidebar_animation.setEndValue(target_width)
        self.sidebar_animation.start()

        # Update button text and save state
        if target_width == 0:
            self.toggle_sidebar_button.setText("▶")  # Right arrow to expand
            self.settings.setValue("sidebar_visible", False)
        else:
            self.toggle_sidebar_button.setText("◀")  # Left arrow to collapse
            self.settings.setValue("sidebar_visible", True)

    @Slot()
    def _on_result_selection_changed(self):
        """Load preview when a result is selected."""
        selected_items = self.results_table.selectedIndexes()
        if not selected_items:
            self.preview_pane.clear()
            return

        # Get the file path from the selected row (column 1 is File Path)
        row = selected_items[0].row()
        file_name_item = self.results_table.item(row, 0)
        folder_path_item = self.results_table.item(row, 1)
        if not file_name_item or not folder_path_item:
            return

        file_path = Path(folder_path_item.text()) / file_name_item.text()
        self.preview_pane.handle_preview(str(file_path))

    @Slot()
    def _add_folder(self):
        """Opens a dialog to select a folder and adds it to the list."""
        folder_path = QFileDialog.getExistingDirectory(self, "Select Folder to Index")
        if folder_path and folder_path not in self.searcher.folder_paths:
            self.searcher.folder_paths.add(folder_path)
            self.settings_tab.folder_list.addItem(folder_path)
            self._save_config()  # Saving config will also save the index

    @Slot()
    def _remove_folder(self):
        """Removes the selected folder from the list."""
        selected_items = self.settings_tab.folder_list.selectedItems()
        if not selected_items:
            return
        reply = QMessageBox.question(
            self,
            "Remove folder from monitoring",
            "Confirm removing the selected folder(s)? This will trigger index rebuilding.",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
        )
        if reply == QMessageBox.StandardButton.Yes:
            for item in selected_items:
                self.searcher.folder_paths.discard(item.text())
                self.settings_tab.folder_list.takeItem(self.settings_tab.folder_list.row(item))
            self._save_config()  # Save config will also save the index

    @Slot()
    def _add_file(self):
        """Opens a dialog to select a file and adds it to the list."""
        file_path, _ = QFileDialog.getOpenFileName(self, "Select File to Index", "")
        if file_path and file_path not in self.searcher.file_paths:
            self.searcher.file_paths.add(file_path)
            self.settings_tab.file_list.addItem(file_path)
            self._save_config()  # Save config will also save the index

    @Slot()
    def _remove_file(self):
        """Removes the selected file from the list."""
        selected_items = self.settings_tab.file_list.selectedItems()
        if not selected_items:
            return
        reply = QMessageBox.question(
            self,
            "Remove file from monitoring",
            "Confirm removing the selected file(s)? This will trigger index rebuilding.",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
        )
        if reply == QMessageBox.StandardButton.Yes:
            for item in selected_items:
                self.searcher.file_paths.discard(item.text())
                self.settings_tab.file_list.takeItem(self.settings_tab.file_list.row(item))
            self._save_config()  # Save config will also save the index

    @Slot()
    def _start_rebuild_index(self):
        """Starts the index rebuilding process in a separate thread."""
        if self.index_rebuilder_thread.isRunning():
            # self.status_bar.showMessage("Index rebuilding is already in progress. Please wait.")
            logger.info(f"Index rebuilding triggered while the rebuilding thread is already running.")
            return
        self.status_bar.showMessage("Starting to rebuild index...")
        self.settings_tab.rebuild_index_button.setEnabled(False)
        self.search_button.setEnabled(False)
        self.index_rebuilder_thread.start()

    @Slot()
    def _on_rebuild_started(self):
        """Updates UI when index rebuilding starts."""
        self.status_bar.showMessage("Index is being rebuilt... Might take a while.")
        logger.info("Rebuild started signal received.")
        self.rebuild_start_time = time.time()

    @Slot(bool)
    def _on_rebuild_finished(self, success):
        """Updates UI when index rebuilding finishes."""
        if success:
            time_elapsed = time.time() - self.rebuild_start_time
            self.status_bar.showMessage(
                f"Index rebuild finished in {time_elapsed:.1f} seconds"
            )
            logger.info(
                f"Rebuild finished signal received: Success. Time used {time_elapsed} seconds"
            )

        else:
            self.status_bar.showMessage("Index rebuild failed.")
            logger.info("Rebuild finished signal received: Failure.")
        self.settings_tab.rebuild_index_button.setEnabled(True)
        self.search_button.setEnabled(True)

        # Ensure file monitor is always running and up-to-date after an index operation (load or rebuild)
        if self.file_monitor_thread and self.file_monitor_thread.isRunning():
            self.file_monitor_thread.stop()  # Stop if already running to restart with fresh mtimes
        # Re-initialize the thread to ensure it picks up latest paths and then start it
        interval = self.settings.value("monitor_interval", FILE_CHANGE_CHECK_INTERVAL, type=int)
        self.file_monitor_thread = FileChangeMonitorThread(
            self.searcher, interval=interval
        )
        self.file_monitor_thread.changes_detected.connect(
            self._on_file_changes_detected
        )
        self.file_monitor_thread.monitoring_status_update.connect(
            self.status_bar.showMessage
        )
        self.file_monitor_thread.start()

    @Slot()
    def _on_file_changes_detected(self):
        logger.info("File changes detected! Triggering index rebuild.")
        self.status_bar.showMessage(
            "File changes detected! Triggering index rebuild..."
        )
        self._start_rebuild_index()

    @Slot()
    def _perform_search(self):
        """Performs a search based on the current query and selected mode."""
        query = self.query_input.text().strip()
        if not query:
            QMessageBox.warning(self, "Empty Query", "Please enter a search query.")
            return

        if self.index_rebuilder_thread.isRunning():
            QMessageBox.information(
                self,
                "Indexing in Progress",
                "Cannot search while index is being rebuilt. Please wait.",
            )
            return

        selected_mode = SEARCH_MODE_MAPPING[self.search_mode_combo.currentText()]
        logger.info(f"Searching for '{query}' with mode: {selected_mode}")
        self.status_bar.showMessage(f"Searching for '{query}'...")

        dedup = self.settings.value("dedup_checkbox_checked", False, type=bool)
        results, highlight_tokens = self.searcher.search(query, mode=selected_mode,deduplicate=dedup)
        # store raw results so filters can be applied locally without triggering new searches
        self.last_search_results = results
        self.last_highlight_tokens = highlight_tokens
        self.last_query = query
        self._update_filter_counts(self.last_search_results)
        filtered_results = self._apply_filters(self.last_search_results)
        self._display_results(
            filtered_results, self.last_highlight_tokens, self.last_query
        )
        self.status_bar.showMessage(
            f"Search completed. Found {len(filtered_results)} results"
        )
        self._save_config()  # Save the last used search mode

    def _apply_filters(self, results):
        """Applies selected filters to search results."""
        filtered_results = results

        # Apply date filters
        selected_date_filter = None
        for label, (checkbox, delta) in self.date_filters.items():
            if checkbox.isChecked() and label != "All Time":
                selected_date_filter = delta
                break

        if selected_date_filter:
            cutoff_time = datetime.now() - selected_date_filter
            filtered_results = [
                result
                for result in filtered_results
                if result.get("mtime")
                   and datetime.fromtimestamp(result["mtime"]) >= cutoff_time
            ]

        # Apply file type filters
        selected_file_types = []
        for label, (checkbox, extensions) in self.file_type_filters.items():
            if checkbox.isChecked() and label != "All Types":
                if extensions:
                    selected_file_types.extend(extensions)

        if selected_file_types:
            filtered_results = [
                result
                for result in filtered_results
                if any(
                    result.get("path", "").lower().endswith(ext)
                    for ext in selected_file_types
                )
            ]

        return filtered_results

    @Slot()
    def _on_filter_changed(self):
        """Keep "All" options exclusive with specifics, then refresh counts/results.

        Prioritize direct clicks on the "All" checkbox so it immediately takes effect
        (instead of being overridden by the presence of previously-checked specifics).
        """
        sender = self.sender()

        def enforce_exclusive(filters, all_key):
            all_cb = filters[all_key][0]
            # If the user clicked the "All" checkbox, honor that action first
            if sender is all_cb:
                if all_cb.isChecked():
                    # Uncheck all specific options
                    for k, (cb, _) in filters.items():
                        if k != all_key:
                            cb.blockSignals(True)
                            cb.setChecked(False)
                            cb.blockSignals(False)
                else:
                    # If user unchecked All and no specific is selected, keep All checked
                    any_specific = any(
                        cb.isChecked() for k, (cb, _) in filters.items() if k != all_key
                    )
                    if not any_specific:
                        all_cb.blockSignals(True)
                        all_cb.setChecked(True)
                        all_cb.blockSignals(False)
            else:
                # A specific checkbox changed: if any specific is checked -> uncheck All, else check All
                specific_checked = any(
                    cb.isChecked() for k, (cb, _) in filters.items() if k != all_key
                )
                all_cb.blockSignals(True)
                all_cb.setChecked(not specific_checked)
                all_cb.blockSignals(False)

        enforce_exclusive(self.date_filters, "All Time")
        enforce_exclusive(self.file_type_filters, "All Types")

        # Reapply filters using the last stored search results (no new search)
        if self.last_search_results is not None:
            self._update_filter_counts(self.last_search_results)
            filtered = self._apply_filters(self.last_search_results)
            self._display_results(filtered, self.last_highlight_tokens, self.last_query)

    @Slot()
    def _clear_filters(self):
        """Clears all filters and resets to default."""
        # Reset date filters to "All Time"
        for label, (checkbox, _) in self.date_filters.items():
            checkbox.blockSignals(True)
            checkbox.setChecked(label == "All Time")
            checkbox.blockSignals(False)

        # Reset file type filters to "All Types"
        for label, (checkbox, _) in self.file_type_filters.items():
            checkbox.blockSignals(True)
            checkbox.setChecked(label == "All Types")
            checkbox.blockSignals(False)

        # Reapply filters using stored last search results (do not trigger a new search)
        if self.last_search_results is not None:
            self._update_filter_counts(self.last_search_results)
            filtered = self._apply_filters(self.last_search_results)
            self._display_results(filtered, self.last_highlight_tokens, self.last_query)

    def _update_filter_counts(self, results):
        """Updates the counts displayed next to filter checkboxes based on search results."""
        # Update date filter counts
        for label, (checkbox, delta) in self.date_filters.items():
            if label == "All Time":
                count = len(results)
            else:
                cutoff_time = datetime.now() - delta
                count = len(
                    [
                        r
                        for r in results
                        if r.get("mtime")
                           and datetime.fromtimestamp(r["mtime"]) >= cutoff_time
                    ]
                )
            checkbox.setText(f"{label} ({count})")

        # Update file type filter counts
        for label, (checkbox, extensions) in self.file_type_filters.items():
            if label == "All Types":
                count = len(results)
            else:
                if extensions:
                    count = len(
                        [
                            r
                            for r in results
                            if any(
                            r.get("path", "").lower().endswith(ext)
                            for ext in extensions
                        )
                        ]
                    )
                else:
                    count = 0
            checkbox.setText(f"{label} ({count})")

    def _display_results(self, results, highlight_tokens=None, query=""):
        """Displays search results in the QTableWidget. query_tokens used for highlighting"""
        self.results_table.setRowCount(0)  # Clear existing results
        highlight_words = highlight_tokens

        for row_idx, result in enumerate(results):
            self.results_table.insertRow(row_idx)
            file_name = os.path.basename(result.get("path", "N/A"))
            file_path = result.get("path", "N/A")
            content = result.get("content", "N/A")
            file_size = (
                f"{result.get('size', 0) / 1024:.2f} KB"
                if result.get("size") is not None
                else "N/A"
            )
            mtime_timestamp = result.get("mtime")
            date_modified = (
                datetime.fromtimestamp(mtime_timestamp).strftime("%Y-%m-%d %H:%M:%S")
                if mtime_timestamp
                else "N/A"
            )
            # source_mode = result.get('source_mode', 'N/A')

            _item = QTableWidgetItem(file_name)
            _item.setToolTip(file_name)
            self.results_table.setItem(row_idx, 0, _item)
            folder_path = os.path.dirname(file_path)
            # making the forward and backward slashes consistent
            folder_path = str(Path(folder_path))
            _item = QTableWidgetItem(folder_path)
            _item.setToolTip(folder_path)
            self.results_table.setItem(row_idx, 1, _item)
            # self.results_table.setCellWidget(row_idx, 2, label)
            # item.setTextAlignment(Qt.AlignLeft | Qt.AlignTop)

            # text_edit1.setReadOnly(True)
            text_edit = QTextEdit()
            text_edit.setHtml(highlight_keywords_in_text(content, highlight_words, query))
            text_edit.setMaximumHeight(80)
            # text_edit.setStyleSheet('')
            # text_edit.setStyle(QApplication.style())
            self.results_table.setCellWidget(row_idx, 2, text_edit)

            # self.results_table.setItem(row_idx, 2, QTableWidgetItem(snippet))
            self.results_table.setItem(row_idx, 3, QTableWidgetItem(file_size))
            self.results_table.setItem(row_idx, 4, QTableWidgetItem(date_modified))
            # self.results_table.setItem(row_idx, 5, QTableWidgetItem(source_mode))

        self.results_table.setWordWrap(True)
        self.results_table.resizeRowsToContents()  # Adjust row heights to fit content

        # self.results_table.resizeColumnsToContents()  # Adjust column widths to fit content

    def on_tray_icon_activated(self, reason):
        if reason == QSystemTrayIcon.Trigger:  # Left click
            self.show()
            self.raise_()
            self.activateWindow()

    @Slot(str)
    def on_rebuild_signal_progress(self, msg):
        self.status_bar.showMessage(msg)

    def closeEvent(self, event):
        """Handles the window close event."""
        if self.settings.value("minimize_to_tray", True, type=bool):
            event.ignore()
            self.hide()
            # self.tray_icon.showMessage("Tray Application", "Application minimized to tray.",
            #                            QSystemTrayIcon.Information, 2000)
        else:
            self.tray_icon.hide()
            logger.info("Application closing. Stopping threads...")
            if self.file_monitor_thread and self.file_monitor_thread.isRunning():
                logger.info("Stopping file monitor thread")
                self.file_monitor_thread.stop()
            if self.index_rebuilder_thread and self.index_rebuilder_thread.isRunning():
                # self.index_rebuilder_thread.wait()
                logger.info("Stopping rebuilder thread")
                self.index_rebuilder_thread.stop()
            logger.info("Saving config")
            self._save_config()
            self.settings.setValue("window_geometry", self.saveGeometry())
            header = self.results_table.horizontalHeader()
            self.settings.setValue("headerState", header.saveState())
            self.settings.setValue("splitterState", self.splitter.saveState())
            self.settings.setValue("font_size", int(self.settings_tab.spin_box.value()))
            logger.info("Config saved")
            log_file.close()
            super().closeEvent(event)

    def close_app(self):
        self.tray_icon.hide()
        QApplication.instance().quit()


class CustomSplashScreen(QWidget):
    def __init__(
            self,
            size=(420, 260),
            title="Local Search",
            subtitle="Loading inference engine and models...",
            logo_path=str(SCRIPT_DIR / "resources/icon.png"),
            background_color="#ffffff",
            radius=16,
    ):
        self._size = size
        # non-zero radius causes black corners for pyinstaller exe
        if is_running_in_pyinstaller():
            self._radius = 0
        else:
            self._radius = radius
        self._background_color = QColor(background_color)

        super().__init__(None,Qt.FramelessWindowHint)
        # super().__init__(None,
        #                  Qt.FramelessWindowHint | Qt.WindowStaysOnTopHint)
        self.setFixedSize(self._size[0], self._size[1])

        self.setAttribute(Qt.WA_TranslucentBackground, True)
        self.setAttribute(Qt.WA_NoSystemBackground, True)
        self.setAttribute(Qt.WA_OpaquePaintEvent, True)
        self.setWindowFlag(Qt.FramelessWindowHint)

        self.title = title
        self.subtitle = subtitle
        self.logo = QPixmap(logo_path) if logo_path else None

        self._opacity_anim = None

        self._center_on_screen()

    def paintEvent(self, event):
        painter = QPainter(self)
        painter.setRenderHint(QPainter.Antialiasing)

        rect = self.rect()

        # Background
        painter.setBrush(self._background_color)
        painter.setPen(Qt.NoPen)
        painter.drawRoundedRect(rect, self._radius, self._radius)

        y = 40

        # Logo
        if self.logo:
            logo_size = 100
            scaled_logo = self.logo.scaled(
                logo_size * self.devicePixelRatio(),  # Scale to physical pixels
                logo_size * self.devicePixelRatio(),
                Qt.KeepAspectRatio,
                Qt.SmoothTransformation
            )
            scaled_logo.setDevicePixelRatio(self.devicePixelRatio())
            logo_rect = QRect((rect.width() - logo_size) // 2, y, logo_size, logo_size)
            painter.drawPixmap(logo_rect, scaled_logo)
            y += logo_size + 12

        # Title
        painter.setPen(QColor("#2c3e50"))
        painter.setFont(QFont("Segoe UI", 18, QFont.DemiBold))
        painter.drawText(
            QRect(0, y, rect.width(), 32),
            Qt.AlignCenter,
            self.title,
        )
        y += 32

        # Subtitle
        painter.setFont(QFont("Segoe UI", 11))
        painter.setPen(QColor("#808080"))
        painter.drawText(
            QRect(0, y, rect.width(), 24),
            Qt.AlignCenter,
            self.subtitle,
        )

    def _center_on_screen(self):
        screen = QApplication.primaryScreen().availableGeometry()
        self.move(
            screen.center() - self.rect().center()
        )


def main():
    app = QApplication(sys.argv)
    is_autostart = AUTOSTART_ARG in sys.argv
    use_gpu = '--gpu' in sys.argv
    if use_gpu:
        configs.DEVICE='gpu'
    logger.info(f"Application launched. Is autostart: {is_autostart}")
    app.settings = QSettings("LocalSearch", "LocalSearch")
    app.setWindowIcon(QIcon(str(SCRIPT_DIR / "resources/icon.png")))
    font_id = QFontDatabase.addApplicationFont(str(SCRIPT_DIR / "resources/SourceHanSansSC-Normal.otf"))
    font_families = QFontDatabase.applicationFontFamilies(font_id)
    custom_font_family = font_families[0]
    font_size = app.settings.value("font_size", 12, type=int)
    custom_font = QFont(custom_font_family, font_size)
    app.setFont(custom_font)

    if is_autostart:
        pass
    else:
        # splash screen
        splash = CustomSplashScreen()

        splash.show()
        app.processEvents()

    window = MainWindow()
    from qt_material import apply_stylesheet

    THEME_EXTRA_SETTINGS = {
        "density_scale": "-2",  # Lower density (more compact layout)
        "font_size": font_size,
        "font_family": custom_font_family,
    }
    apply_stylesheet(
        app,
        theme="light_blue.xml",
        invert_secondary=True,
        extra=THEME_EXTRA_SETTINGS,
        css_file=str(SCRIPT_DIR / "resources/custom.css"),
    )
    if is_autostart:
        window.setWindowState(Qt.WindowMinimized)
    else:

        window.show()
        window.raise_()
        window.activateWindow()
        splash.close()


    app.exec_()
    # sys.exit(app.exec())
if __name__ == "__main__":
    main()




