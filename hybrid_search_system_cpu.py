# This is an app that performs hybrid search (BM25 + embedding similarity) on local files and folders.
# It uses FAISS for embedding similarity and rank_bm25 for BM25 search.
# The app has a GUI built with PySide6.


import os, sys, pickle
from pathlib import Path

os.environ["DO_NOT_TRACK"] = "true"  # for speeding up unstructured
os.environ["HF_HUB_OFFLINE"] = "1"
import json
import time
from datetime import datetime, timedelta

from PySide6.QtWidgets import (
    QApplication,
    QMenu,
    QMainWindow,
    QWidget,
    QVBoxLayout,
    QHBoxLayout,
    QPushButton,
    QLineEdit,
    QListWidget,
    QTableWidget,
    QSystemTrayIcon,
    QTableWidgetItem,
    QHeaderView,
    QComboBox,
    QFileDialog,
    QTabWidget,
    QSpinBox,
    QStatusBar,
    QMessageBox,
    QLabel,
    QAbstractItemView,
    QTextEdit,
    QCheckBox,
    QSplashScreen,
    QSizePolicy,
    QScrollArea,
    QSplitter,
    QStackedWidget,
    QTextBrowser,
)
from PySide6.QtCore import (
    QThread,
    Signal,
    Slot,
    Qt,
    QUrl,
    QSettings,
    QByteArray,
    QObject,
    QPropertyAnimation,
    QSize,
)
from PySide6.QtGui import (
    Qt,
    QAction,
    QFont,
    QFontDatabase,
    QDesktopServices,
    QIcon,
    QPixmap,
    QColor,
    QImage,
)

from qt_material import apply_stylesheet, QtStyleTools

import re
from html import escape

import logging

import fitz  # PyMuPDF
import markdown
from docx import Document
from pptx import Presentation
from configs import *

# TODO 处理supported extensions. 添加合理的非preview类型，例如mp3, mp4, . 删除preview涉及到的.json等编程类型
# TODO 解决self.preview_pane.setText 的问题


os.makedirs(INDEX_DIR, exist_ok=True)
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
    handlers=[
        logging.FileHandler(LOG_PATH, encoding="utf-8"),
        logging.StreamHandler(sys.stdout),
    ],
)
logger = logging.getLogger(__name__)


def import_and_initialize():
    """
    Import the time large libraries
    Returns:

    """
    global np, faiss, bm25s, Tokenizer, Stemmer
    global partition, RecursiveCharacterTextSplitter, jieba, PdfReader
    global DEVICE
    import numpy as np

    logger.info("Loading faiss")
    import faiss
    import bm25s
    from bm25s.tokenization import Tokenizer
    import Stemmer
    from tokenizers import Tokenizer as EmbTokenizer

    logger.info("Loading Unstructured")
    from unstructured.partition.auto import partition
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from pypdf import PdfReader
    import rjieba as jieba

    logger.info("Loading finished")

    global EMBEDDING_TOKENIZER
    DEVICE = "cpu"
    import onnxruntime as ort

    global session
    EMBEDDING_TOKENIZER = EmbTokenizer.from_file(EMBEDDING_TOKENIZER_FILE)
    EMBEDDING_TOKENIZER.enable_padding(length=EMBEDDING_MODEL_TOKEN_LENGTH)
    session = ort.InferenceSession(
        EMBEDDING_MODEL_ONNX_FILE, providers=["CPUExecutionProvider"]
    )


class HybridSearch(QObject):
    """
    A class for performing hybrid search (BM25 for chunks, BM25 for documents,
    embedding similarity, and a combination of all) on files and folders.
    It uses FAISS for embedding similarity and rank_bm25 for BM25 search.
    """

    index_rebuild_progress_signal = Signal(str)

    def __init__(self, file_paths=None, folder_paths=None):
        """
        Initializes the HybridSearch object.

        Args:
            file_paths (list): A list of initial file paths to index.
            folder_paths (list): A list of initial folder paths to index.
        """
        super().__init__()
        self.file_paths = set(file_paths) if file_paths else set()
        self.folder_paths = set(folder_paths) if folder_paths else set()

        # Data structures to store indexed content and metadata
        self.documents = (
            []
        )  # List of dictionaries: {'path': ..., 'content': ..., 'mtime': ..., 'size': ...}
        self.chunks = []  # List of dictionaries: {'doc_idx': ..., 'text': ...}
        self.old_chunks = None
        self.faiss_index = None
        self.bm25_chunk_model = None
        self.bm25_doc_model = None
        self.embedding_model = None

        self.chunk_size = 200  # Characters per chunk
        self.chunk_overlap = 30  # Overlap between chunks

        self.stemmer = Stemmer.Stemmer("english")
        self.tokenizer_doc = Tokenizer(
            stemmer=self.stemmer, stopwords=STOPWORDS, splitter=self._tokenize
        )
        self.tokenizer_chunk = Tokenizer(
            stemmer=self.stemmer, stopwords=STOPWORDS, splitter=self._tokenize
        )
        self._load_embedding_model()
        self.splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size, chunk_overlap=self.chunk_overlap
        )

    def _load_embedding_model(self):
        self.device = DEVICE

        self.embedding_tokenizer = EMBEDDING_TOKENIZER
        self.embedding_model = session
        logger.info("Embedding model loaded successfully.")

    def _encode(self, texts, batch_size=32):
        all_outputs = []
        for i in range(0, len(texts), batch_size):
            batch_texts = texts[i : i + batch_size]
            res = EMBEDDING_TOKENIZER.encode_batch_fast(batch_texts)
            input_ids = [e.ids for e in res]
            attention_mask = [e.attention_mask for e in res]
            inputs = {"input_ids": input_ids, "attention_mask": attention_mask}

            outputs = session.run(None, inputs)
            all_outputs.append(outputs[1])  # Assuming the first output is needed
        if len(all_outputs) >= 1:
            embeddings = np.concatenate(all_outputs, axis=0)
        else:
            embeddings = []
        return embeddings

    @staticmethod
    def _read_file_content(file_path):
        # avoid bloat in unstructured[pdf]
        logger.info(f"Parsing{file_path}")
        try:
            if str(Path(file_path).suffix) == ".pdf":
                # elements = partition(file_path, strategy='fast')
                reader = PdfReader(file_path)
                full_text = "\n".join([page.extract_text() for page in reader.pages])
            elif str(Path(file_path).suffix) in [
                ".png",
                ".jpg",
                ".jpeg",
                ".bmp",
                ".gif",
                ".tiff",
            ]:
                return ""
            else:
                elements = partition(str(Path(file_path)))
                full_text = "\n".join(
                    element.text for element in elements if element.text
                )
            logger.info(f"Successfully parsed file {file_path}")
            return full_text
        except Exception as e:
            logger.info(f"Can't parse file content for {file_path}, {e}")
            return None

    @staticmethod
    def _tokenize(text):
        return jieba.cut_for_search(text.lower())

    def rebuild_index(self):
        """
        Rebuilds the entire search index from scratch.
        This involves reading files, chunking, generating embeddings,
        and building BM25 and FAISS indexes.
        """
        if not self.embedding_model:
            logger.info("Embedding model not loaded. Cannot rebuild index.")
            return False

        logger.info("Rebuilding index...")
        self.documents = []
        self.old_chunks = self.chunks
        self.chunks = []
        all_chunk_texts = []
        self.faiss_index = None
        self.bm25_chunk_model = None
        self.bm25_doc_model = None
        self.tokenizer_doc = Tokenizer(
            stemmer=self.stemmer, stopwords=STOPWORDS, splitter=self._tokenize
        )
        self.tokenizer_chunk = Tokenizer(
            stemmer=self.stemmer, stopwords=STOPWORDS, splitter=self._tokenize
        )

        files_to_index = set()
        for folder_path in self.folder_paths:
            if os.path.isdir(folder_path):
                for root, _, files in os.walk(folder_path):
                    for file in files:
                        if str(Path(file).suffix) in SUPPORTED_EXTENSIONS:
                            files_to_index.add(os.path.join(root, file))
        files_to_index.update(self.file_paths)
        total_files = len(files_to_index)
        for idx, file_path in enumerate(files_to_index):
            self.index_rebuild_progress_signal.emit(
                f"Rebuild index,{idx + 1}/{total_files}, {str(Path(file_path))}"
            )
            if not os.path.exists(file_path):
                logger.info(f"File not found, skipping: {file_path}")
                continue

            content = self._read_file_content(file_path)
            if content is None:  # reading file contents for images would return ""
                continue

            try:
                mtime = os.path.getmtime(file_path)
                size = os.path.getsize(file_path)
            except OSError as e:
                logger.info(f"Could not get metadata for {file_path}: {e}")
                continue

            doc_idx = idx
            self.documents.append(
                {
                    "path": str(Path(file_path)),
                    "content": content,
                    "mtime": mtime,
                    "size": size,
                    # 'emb_needs_update': emb_needs_update
                }
            )

            file_chunks = self.splitter.split_text(content)
            for chunk_text in file_chunks:
                self.chunks.append(
                    {
                        "doc_idx": doc_idx,
                        "text": chunk_text,
                        # 'emb_needs_update': emb_needs_update
                    }
                )
                all_chunk_texts.append(chunk_text)

        if not self.documents:
            logger.info("No documents to index.")
            return True

        # Build BM25 Document Model
        logger.info("Rebuilding index, setting up tokenizer")
        self.index_rebuild_progress_signal.emit(
            f"Rebuilding index...Processing keyword index..."
        )
        corpus_doc = [doc["path"] + doc["content"] for doc in self.documents]
        corpus_doc_tokens = self.tokenizer_doc.tokenize(
            corpus_doc, update_vocab=True, allow_empty=True
        )
        self.bm25_doc_model = bm25s.BM25(method="bm25+")
        self.bm25_doc_model.index(corpus_doc_tokens)
        logger.info(f"BM25 Document Model built with {len(self.documents)} documents.")

        # Build BM25 Chunk Model
        if all_chunk_texts:
            logger.info("Rebuilding index, preparing to build chunk BM25 model")
            corpus_chunk_tokens = self.tokenizer_chunk.tokenize(
                all_chunk_texts, update_vocab=True, allow_empty=True
            )
            self.bm25_chunk_model = bm25s.BM25(method="bm25+")
            self.bm25_chunk_model.index(corpus_chunk_tokens)
            logger.info(f"BM25 Chunk Model built with {len(self.chunks)} chunks.")
            logger.info(f"Generating embeddings for {len(all_chunk_texts)} chunks...")

            # reuse embeddings in self.old_chunks
            chunks_subset_to_update = []
            for chunk in self.chunks:
                match = next(
                    (
                        c
                        for c in self.old_chunks
                        if c["doc_idx"] == chunk["doc_idx"]
                        and c["text"] == chunk["text"]
                    ),
                    None,
                )
                if match and "embedding" in match.keys():
                    chunk["embedding"] = match["embedding"]
                else:
                    chunks_subset_to_update.append(chunk)
            texts_to_embed = [chunk["text"] for chunk in chunks_subset_to_update]
            self.index_rebuild_progress_signal.emit(
                f"Rebuilding index...Calculating embeddings..."
            )
            chunks_subset_embeddings = self._encode(texts_to_embed)
            for idx, emb in enumerate(chunks_subset_embeddings):
                chunks_subset_to_update[idx]["embedding"] = emb

            chunk_embeddings = np.stack(
                [chunk["embedding"] for chunk in self.chunks], axis=0
            )
            d = chunk_embeddings.shape[1]  # Dimension of embeddings
            param = "HNSW64"
            self.faiss_index = faiss.index_factory(d, param, faiss.METRIC_INNER_PRODUCT)
            self.faiss_index.add(chunk_embeddings)
            logger.info(f"FAISS index built with {self.faiss_index.ntotal} embeddings.")

        else:
            logger.info("No chunks to index for embeddings and BM25 chunk model.")
        logger.info("Index rebuilding complete.")
        self.index_rebuild_progress_signal.emit(f"Index successfully rebuilt.")
        return True

    def search(self, query, mode="hybrid", top_k=50):
        """
        Performs a search based on the specified mode.

        Args:
            query (str): The search query.
            mode (str): The search mode ('bm25_chunk', 'bm25_document', 'embedding', 'hybrid').
            top_k (int): The number of top results to retrieve for each sub-search.

        Returns:
            list: A list of dictionaries, each representing a search result.
                  Each dictionary contains file metadata
        """
        assert mode in ["hybrid", "bm25_document", "embedding"]
        results = []
        query_tokens = self.tokenizer_chunk.tokenize([query], allow_empty=False)
        highlight_tokens = self.tokenizer_chunk.tokenize([query], allow_empty=False, return_as="string")[0]  # type: ignore
        if mode == "hybrid":  # bm25_chunk retrieval
            if self.bm25_chunk_model and self.chunks:
                logger.info(f"Performing BM25 chunk search for '{query}'...")
                k = min(top_k, len(self.chunks))
                tmp, scores = self.bm25_chunk_model.retrieve(query_tokens, k=k)
                chunk_scores = scores[0]
                ranked_chunk_indices = np.argsort(chunk_scores)[::-1]

                bm25_chunk_results = []
                for idx in ranked_chunk_indices:
                    if idx < len(self.chunks):  # Ensure index is valid
                        chunk = self.chunks[idx]
                        doc = self.documents[chunk["doc_idx"]]
                        # sometimes the result does not contain any query token.
                        # can't use score to discern
                        if not any(word in chunk["text"] for word in highlight_tokens):
                            continue
                        bm25_chunk_results.append(
                            {
                                "path": doc["path"],
                                "content": chunk["text"],
                                "mtime": doc["mtime"],
                                "size": doc["size"],
                                "score": chunk_scores[idx],
                                "source_mode": "BM25 Chunk",
                            }
                        )
                results.extend(bm25_chunk_results)

        if mode == "bm25_document" or mode == "hybrid":
            if self.bm25_doc_model and self.documents:
                logger.info(f"Performing BM25 document search for '{query}'...")
                # Get top-k results as a tuple of (doc ids, scores). Both are arrays of shape (n_queries, k).
                # To return docs instead of IDs, set the `corpus=corpus` parameter.
                k = min(top_k, len(self.documents))
                tmp, scores = self.bm25_doc_model.retrieve(query_tokens, k=k)
                doc_scores = scores[0]
                ranked_doc_indices = np.argsort(doc_scores)[::-1]

                bm25_doc_results = []
                for idx in ranked_doc_indices:
                    if idx < len(self.documents):
                        doc = self.documents[idx]
                        if not any(word in doc["content"] for word in highlight_tokens):
                            continue
                        bm25_doc_results.append(
                            {
                                "path": doc["path"],
                                "content": doc["content"],
                                "mtime": doc["mtime"],
                                "size": doc["size"],
                                "score": doc_scores[idx],
                                "source_mode": "BM25 Document",
                            }
                        )
                bm25_doc_results = sorted(
                    bm25_doc_results, key=lambda d: d["score"], reverse=True
                )
                results.extend(bm25_doc_results)

        if mode == "embedding" or mode == "hybrid":
            if self.faiss_index and self.embedding_model:
                logger.info(f"Performing embedding similarity search for '{query}'...")
                query_embedding = self._encode([query])
                D, I = self.faiss_index.search(
                    query_embedding, top_k
                )  # D: distances, I: indices

                embedding_results = []
                for i, score in zip(I[0], D[0]):
                    if i != -1 and i < len(self.chunks):  # Check for valid index
                        chunk = self.chunks[i]
                        doc = self.documents[chunk["doc_idx"]]
                        embedding_results.append(
                            {
                                "path": doc["path"],
                                "content": chunk["text"],
                                "mtime": doc["mtime"],
                                "size": doc["size"],
                                "score": -score,  # FAISS returns distance, convert to similarity score
                                "source_mode": "Embedding",
                            }
                        )
                embedding_results = sorted(
                    embedding_results, key=lambda d: d["score"], reverse=True
                )
                results.extend(embedding_results)

        # For hybrid mode, combine, and use reciprocal rank fusion
        if mode == "hybrid":
            for i, item in enumerate(bm25_doc_results):
                item["rank"] = i + 1
            for i, item in enumerate(embedding_results):
                item["rank"] = i + 1
            for i, item in enumerate(bm25_chunk_results):
                item["rank"] = i + 1
            k = 60  # parameter for rrf
            for res in results:
                r_b_chunk = (
                    res["rank"]
                    if res["source_mode"] == "BM25 Chunk"
                    else len(bm25_chunk_results) + 1
                )
                r_b_doc = (
                    res["rank"]
                    if res["source_mode"] == "BM25 Document"
                    else len(bm25_doc_results) + 1
                )
                r_e = (
                    res["rank"]
                    if res["source_mode"] == "Embedding"
                    else len(embedding_results) + 1
                )
                rrf_score = 1 / (k + r_b_chunk) + 1 / (k + r_b_doc) + 1 / (k + r_e)
                res["rrf_score"] = rrf_score
                # count how many highlight_tokens appeared in the content
                res["highlight_tokens_appearance"] = sum(
                    1 for ht in highlight_tokens if ht in res["content"]
                )
            results = sorted(results, key=lambda x: x["rrf_score"], reverse=True)

            # stage 2 ranking. prioritize results containing more hits for query
            best_hits = []
            good_hits = []
            n_highlight_tokens = len(highlight_tokens)
            for i, res in enumerate(results[:]):  # iterate over a copy
                if query in res["content"]:
                    best_hits.append(res)
                    results.remove(res)
                elif (
                    n_highlight_tokens > 2
                    and res["highlight_tokens_appearance"] > 0.8 * n_highlight_tokens
                ):
                    good_hits.append(res)
                    results.remove(res)
            tmp_results = best_hits + good_hits + results
            # deduplicate based on both content and file path (composite key)
            final_results = []
            seen_pairs = set()
            for res in tmp_results:
                if res["source_mode"] in ["BM25 Chunk", "Embedding"]:
                    # Create a composite key from content and path
                    pair = (res["content"], res["path"])
                    if pair not in seen_pairs:
                        final_results.append(res)
                        seen_pairs.add(pair)
            results = final_results
        return results[:top_k], highlight_tokens

    def save_index_components(self, index_dir=INDEX_DIR):
        os.makedirs(index_dir, exist_ok=True)
        docs_chunks_data = {
            "documents": self.documents,
            "chunks": self.chunks,
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap,
        }
        # skip saving if the data have not changed
        old_docs_chunks_data = {
            "documents": self.documents,
            "chunks": self.old_chunks,
            "chunk_size": self.chunk_size,
            "chunk_overlap": self.chunk_overlap,
        }
        if self.faiss_index is not None:
            faiss.write_index(
                self.faiss_index, os.path.join(index_dir, "faiss_index.bin")
            )

        self.tokenizer_doc.save_vocab(save_dir=str(Path(index_dir) / "tokenizer_doc"))
        self.tokenizer_chunk.save_vocab(
            save_dir=str(Path(index_dir) / "tokenizer_chunk")
        )
        if docs_chunks_data != old_docs_chunks_data:
            with open(os.path.join(index_dir, "docs_chunks_data.pkl"), "wb") as f:
                pickle.dump(docs_chunks_data, f)
        if self.bm25_doc_model and self.bm25_chunk_model:
            self.bm25_doc_model.save(str(Path(index_dir) / "bm25_doc_model"))
            self.bm25_chunk_model.save(str(Path(index_dir) / "bm25_chunk_model"))
        logger.info(f"Index components saved to {index_dir}")
        return True

    def load_index_components(self, index_dir=INDEX_DIR):
        """Loads the index components from disk."""
        if not os.path.exists(index_dir):
            logger.info(f"Index directory not found: {index_dir}")
            return False

        try:
            with open(os.path.join(index_dir, "docs_chunks_data.pkl"), "rb") as f:
                docs_chunks_data = pickle.load(f)
            self.tokenizer_doc.load_vocab(str(Path(index_dir) / "tokenizer_doc"))
            self.tokenizer_chunk.load_vocab(str(Path(index_dir) / "tokenizer_chunk"))

            self.bm25_doc_model = bm25s.BM25.load(
                str(Path(index_dir) / "bm25_doc_model")
            )
            self.bm25_chunk_model = bm25s.BM25.load(
                str(Path(index_dir) / "bm25_chunk_model")
            )

            self.documents = docs_chunks_data.get("documents", [])
            self.chunks = docs_chunks_data.get("chunks", [])
            self.chunk_size = docs_chunks_data.get("chunk_size", 200)
            self.chunk_overlap = docs_chunks_data.get("chunk_overlap", 30)

            faiss_index_path = os.path.join(index_dir, "faiss_index.bin")
            if os.path.exists(faiss_index_path):
                self.faiss_index = faiss.read_index(faiss_index_path)
            else:
                self.faiss_index = None

            logger.info(f"Index components loaded from {index_dir}")
            return True
        except Exception as e:
            logger.info(f"Error loading index components from {index_dir}: {e}")
            # Reset all index related attributes on failure
            self.documents = []
            self.chunks = []
            self.faiss_index = None
            self.bm25_chunk_model = None
            self.bm25_doc_model = None
            return False


class IndexRebuilderThread(QThread):
    """
    A QThread for rebuilding the search index in the background.
    Emits signals to update the UI status.
    """

    rebuild_started = Signal()
    rebuild_finished = Signal(bool)  # True for success, False for failure

    def __init__(self, hybrid_search_instance):
        super().__init__()
        self.hybrid_search = hybrid_search_instance

    def run(self):
        """
        The main execution method for the thread.
        Calls the rebuild_index method of the HybridSearch instance.
        """
        self.rebuild_started.emit()
        success = self.hybrid_search.rebuild_index()
        if success:
            logger.info(
                "Index rebuild successful, now saving index components in background thread..."
            )
            save_success = self.hybrid_search.save_index_components()  # Call save here
            if not save_success:
                logger.info("Warning: Index components failed to save after rebuild.")
                success = False
        self.rebuild_finished.emit(success)

    def stop(self):
        self.terminate()


class FileMonitorThread(QThread):
    """
    A QThread for continuously monitoring indexed files and folders for changes.
    Emits a signal if changes are detected.
    """

    changes_detected = Signal()
    # Signal to update the status bar with monitoring status
    monitoring_status_update = Signal(str)

    def __init__(self, hybrid_search_instance, interval=120):
        super().__init__()
        self.hybrid_search = hybrid_search_instance
        self.interval = interval  # Check interval in seconds
        self._running = True
        self.last_mtimes = self._get_current_mtimes()  # Stores {file_path: mtime}

    def _get_current_mtimes(self):
        """
        Gets the current modification times for all indexed files.
        """
        current_mtimes = {}
        files_to_check = set()

        for folder_path in self.hybrid_search.folder_paths:
            if os.path.isdir(folder_path):
                for root, _, files in os.walk(folder_path):
                    for file in files:
                        if str(Path(file).suffix) in SUPPORTED_EXTENSIONS:
                            files_to_check.add(os.path.join(root, file))
        files_to_check.update(self.hybrid_search.file_paths)

        for file_path in files_to_check:
            try:
                if os.path.exists(file_path):
                    current_mtimes[file_path] = os.path.getmtime(file_path)
            except OSError as e:
                logger.info(f"Error getting mtime for {file_path}: {e}")
        return current_mtimes

    def run(self):
        """
        Continuously monitors files for changes.
        """
        self.last_mtimes = self._get_current_mtimes()
        self.monitoring_status_update.emit("Monitoring file changes...")

        while self._running:
            time.sleep(self.interval)
            logger.info(f"Monitoring check")
            current_mtimes = self._get_current_mtimes()
            changed = False
            # Check for modified or new files
            for file_path, mtime in current_mtimes.items():
                if (
                    file_path not in self.last_mtimes
                    or self.last_mtimes[file_path] != mtime
                ):
                    logger.info(f"Change detected in: {file_path}")
                    changed = True
                    break

            # Check for deleted files/folders
            if not changed:
                for file_path in self.last_mtimes:
                    if file_path not in current_mtimes:
                        logger.info(f"File deleted: {file_path}")
                        changed = True
                        break

            if changed:
                self.changes_detected.emit()
                self.last_mtimes = current_mtimes  # Update after emitting signal

    def stop(self):
        """Stops the monitoring thread."""
        self._running = False
        # self.wait()
        self.terminate()


def create_highlighted_text(text, highlight_words, query, max_length=200):
    # Handle None or empty highlight_words
    if not highlight_words:
        highlight_words = []
    # find the first occurrence of the highlight_words
    perfect_hit_position = text.find(query)  # -1 if not found
    positions_for_words = [
        text.find(word) for word in highlight_words if text.find(word) != -1
    ]
    if perfect_hit_position != -1:
        min_pos = perfect_hit_position
    else:
        min_pos = min(positions_for_words) if positions_for_words else -1
    min_pos = 0 if min_pos - 10 < 0 else min_pos - 10
    truncated_text = text[min_pos : min_pos + max_length] + "..."
    # Escape HTML special characters
    escaped_text = escape(truncated_text)

    if perfect_hit_position != -1:
        pattern = re.compile(re.escape(query), re.IGNORECASE)
        # escaped_text = pattern.sub(lambda m: f"<b>{m.group(0)}</b>", escaped_text)
        escaped_text = pattern.sub(
            lambda m: f"<span style='color:red;'>{m.group(0)}</span>", escaped_text
        )
    else:
        # Highlight specified words in red
        for word in highlight_words:
            pattern = re.compile(re.escape(word), re.IGNORECASE)
            # escaped_text = pattern.sub(lambda m: f"<b>{m.group(0)}</b>", escaped_text)
            escaped_text = pattern.sub(
                lambda m: f"<span style='color:red;'>{m.group(0)}</span>", escaped_text
            )

    return escaped_text


class UniversalPreviewPane(QStackedWidget):
    def __init__(self, parent=None):
        super().__init__(parent)

        # 1. Image & PDF Viewer (via Pixmap)
        self.image_viewer = QLabel()
        self.image_viewer.setAlignment(Qt.AlignCenter)
        self.image_scroll = QScrollArea()
        self.image_scroll.setWidgetResizable(True)
        self.image_scroll.setWidget(self.image_viewer)

        # 2. Markdown & Rich Text Viewer
        self.text_browser = QTextBrowser()
        self.text_browser.setOpenExternalLinks(True)

        # 3. Plain Text / Code Viewer
        self.code_viewer = QTextEdit()
        self.code_viewer.setReadOnly(True)
        self.code_viewer.setLineWrapMode(QTextEdit.NoWrap)

        # 4. Error/Fallback View
        self.error_label = QLabel("Preview not available for this file type.")
        self.error_label.setAlignment(Qt.AlignCenter)

        # Add all to the stack
        self.addWidget(self.image_scroll)  # Index 0
        self.addWidget(self.text_browser)  # Index 1
        self.addWidget(self.code_viewer)  # Index 2
        self.addWidget(self.error_label)  # Index 3

    def handle_preview(self, file_path):
        """Main entry point to update the preview based on file path."""
        if not os.path.exists(file_path):
            self.show_error("File not found.")
            return
        if os.path.getsize(file_path) > 15 * 1024 * 1024:  # multi MB limit
            self.show_error("File too large to preview.")
            return
        ext = os.path.splitext(file_path)[1].lower()

        try:
            if ext in [".jpg", ".jpeg", ".png", ".bmp", ".gif"]:
                self.preview_image(file_path)
            elif ext == ".pdf":
                self.preview_pdf(file_path)
            elif ext == ".md":
                self.preview_markdown(file_path)
            elif ext == ".docx":
                self.preview_docx(file_path)
            elif ext == ".pptx":
                self.preview_pptx(file_path)
            elif ext in [".txt", ".py", ".json", ".csv", ".log", ".xml"]:
                self.preview_text(file_path)
            else:
                self.show_error(f"No preview available for {ext}")
        except Exception as e:
            self.show_error(f"Error loading preview: {str(e)}")

    def preview_image(self, path):
        pixmap = QPixmap(path)
        self._set_pixmap_scaled(pixmap)
        self.setCurrentIndex(0)

    def preview_pdf(self, path):
        """Renders the first page of a PDF as an image."""
        doc = fitz.open(path)
        page = doc.load_page(0)
        pix = page.get_pixmap()
        fmt = QImage.Format_RGBA8888 if pix.alpha else QImage.Format_RGB888
        qimg = QImage(pix.samples, pix.width, pix.height, pix.stride, fmt)
        self._set_pixmap_scaled(QPixmap.fromImage(qimg))
        doc.close()
        self.setCurrentIndex(0)

    def preview_markdown(self, path):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            html = markdown.markdown(f.read())
            # Add basic CSS to make it look nice
            styled_html = f"<style>body {{ font-family: sans-serif; padding: 10px; }}</style>{html}"
            self.text_browser.setHtml(styled_html)
        self.setCurrentIndex(1)

    def preview_docx(self, path):
        doc = Document(path)
        text = [p.text for p in doc.paragraphs[:50]]  # Limit to first 50 paragraphs
        self.text_browser.setPlainText("\n".join(text))
        self.setCurrentIndex(1)

    def preview_pptx(self, path):
        prs = Presentation(path)
        text = []
        for i, slide in enumerate(prs.slides[:5]):  # Limit to first 5 slides
            text.append(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        self.text_browser.setPlainText("\n".join(text))
        self.setCurrentIndex(1)

    def preview_text(self, path):
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            self.code_viewer.setPlainText(f.read(10000))  # Load first 10k chars
        self.setCurrentIndex(2)

    def show_error(self, message):
        self.error_label.setText(message)
        self.setCurrentIndex(3)

    def clear(self):
        self.show_error("")

    def _set_pixmap_scaled(self, pixmap):
        """Helper to scale images to fit the scroll area width."""
        if not pixmap.isNull():
            # Adjust scaling logic based on your preference
            scaled = pixmap.scaled(
                self.size(), Qt.KeepAspectRatio, Qt.SmoothTransformation
            )
            self.image_viewer.setPixmap(scaled)


class SettingsTab(QWidget):
    def __init__(self):
        super().__init__()
        self.settings = QApplication.instance().settings
        self.autostart_checkbox = QCheckBox("Start the application on system startup")
        self.autostart_checkbox.stateChanged.connect(self.on_autostart_checkbox_changed)
        self.help_button = QPushButton("Help")
        self.help_button.clicked.connect(self.on_help_button_pushed)
        self.help_window = None
        self.tray_checkbox = QCheckBox("Minimize to system tray on close")
        self.tray_checkbox.setChecked(
            self.settings.value("minimize_to_tray", True, type=bool)
        )
        self.tray_checkbox.stateChanged.connect(self.toggle_tray_setting)
        self.spin_box = QSpinBox()
        self.spin_box.setRange(8, 20)
        self.spin_box.setFixedWidth(80)
        self.init_ui()
        self.load_settings()

    def init_ui(self):
        layout = QVBoxLayout()
        layout.addWidget(self.autostart_checkbox)
        layout.addWidget(self.tray_checkbox)
        h_layout = QHBoxLayout()
        input_label = QLabel("Font size (takes effect after restart):")
        h_layout.addWidget(input_label)
        h_layout.addWidget(self.spin_box, alignment=Qt.AlignLeft)  # type: ignore
        h_layout.addStretch()
        layout.addLayout(h_layout)
        layout.addWidget(self.help_button)
        layout.addStretch(1)  # Pushes everything to the top
        self.setLayout(layout)

    def load_settings(self):
        autostart_enabled = self.settings.value("autostart_enabled", False, type=bool)
        self.autostart_checkbox.setChecked(autostart_enabled)
        _font_size = self.settings.value("font_size", 10, type=int)
        self.spin_box.setValue(_font_size)

    def toggle_tray_setting(self, state):
        minimize_to_tray = state == 2
        self.settings.setValue("minimize_to_tray", minimize_to_tray)

    def save_settings(self):
        self.settings.setValue("autostart_enabled", self.autostart_checkbox.isChecked())

    def on_help_button_pushed(self):
        to_remove = [".", "[", "]", "'"]
        supported_types = str(SUPPORTED_EXTENSIONS)
        tmp = supported_types.translate(str.maketrans("", "", "".join(to_remove)))
        if sys.platform.startswith("win"):
            conversion_bat_file_location = str(
                Path(os.path.dirname(os.path.abspath(__file__)))
                / "resources"
                / "convert_doc_ppt_to_docx_pptx.bat"
            )
            help_text = f"""使用说明：
            1.在"索引管理"页面添加需要搜索的文件夹或单个文件，随后点击“立即重建索引”按钮。重建索引耗时较长。
            2.支持的文件类型有：{tmp}(对于pdf文件，仅限不需要执行OCR的版本；不支持旧的.doc,.ppt文件)。
            3.在搜索结果列表中，双击文件名即可直接打开文件，双击文件路径可直接打开文件所在路径。
            4.如果搜索结果中的文件名、路径较长导致显示不全，可以单击该单元格后，再次单击即可完整显示。
            5.可以使用脚本文件{conversion_bat_file_location}将指定文件夹下的所有.doc, .ppt文件转换为可检索的.docx, .pptx文件（不会删除旧文件，需要电脑安装了Office）
            6.程序log文件位于{PROGRAM_DATA_PATH}"""
        else:
            help_text = f"""使用说明：
            1.在"索引管理"页面添加需要搜索的文件夹或单个文件，随后点击“立即重建索引”按钮。重建索引耗时较长。
            2.支持的文件类型有：{tmp}(对于pdf文件，仅限不需要执行OCR的版本；不支持旧的.doc,.ppt文件)。
            3.在搜索结果列表中，双击文件名即可直接打开文件，双击文件路径可直接打开文件所在路径。
            4.如果搜索结果中的文件名、路径较长导致显示不全，可以单击该单元格后，再次单击即可完整显示。
            5.程序log文件位于{PROGRAM_DATA_PATH}"""
        help_text = "\n".join(
            line.lstrip() for line in help_text.splitlines()
        )  # remove indentation
        QMessageBox.information(self, "帮助", help_text)

    def on_autostart_checkbox_changed(self, state):
        is_checked = state == Qt.CheckState.Checked.value
        self.save_settings()  # Save the new state immediately

        if is_checked:
            logger.info("Autostart checkbox checked. Attempting to enable autostart...")
            self.enable_autostart()
        else:
            logger.info(
                "Autostart checkbox unchecked. Attempting to disable autostart..."
            )
            self.disable_autostart()

    def enable_autostart(self):
        if sys.platform.startswith("win"):
            import winreg

            try:
                key = winreg.OpenKey(
                    winreg.HKEY_CURRENT_USER,
                    r"Software\Microsoft\Windows\CurrentVersion\Run",
                    0,
                    winreg.KEY_SET_VALUE,
                )
                # Get the path to your executable. This might need to be more robust.
                # For a frozen executable, sys.executable points to the executable.
                # For a script, you might need to run python and the script.
                app_path = f'"{sys.executable}"'  # Or path to your .exe if frozen
                winreg.SetValueEx(
                    key, "MyProgram", 0, winreg.REG_SZ, app_path + " --autostart"
                )
                winreg.CloseKey(key)
                logger.info("Autostart enabled in Windows Registry.")
            except Exception as e:
                QMessageBox.warning(
                    self,
                    "Autostart Error",
                    f"Failed to enable autostart on Windows: {e}",
                )
                logger.info(f"Windows autostart error: {e}")

        elif sys.platform.startswith("linux"):
            # Linux: Create a .desktop file in ~/.config/autostart/
            autostart_dir = os.path.expanduser("~/.config/autostart")
            os.makedirs(autostart_dir, exist_ok=True)
            desktop_file_path = os.path.join(autostart_dir, "myprogram.desktop")
            # Get the path to your executable.
            # For a frozen app, sys.executable is the app path.
            # For a script, you'd need the path to the script and python executable.
            exec_path = (
                sys.executable
            )  # Or the path to your main script, e.g., f"python3 {os.path.abspath(__file__)}"

            desktop_content = f"""[Desktop Entry]
            Type=Application
            Exec=python3 {exec_path} --autostart
            Hidden=false
            NoDisplay=false
            X-GNOME-Autostart-enabled=true
            Name=MyProgram
            Comment=Start MyProgram on boot
            """
            try:
                with open(desktop_file_path, "w") as f:
                    f.write(desktop_content)
                logger.info(f"Autostart .desktop file created at: {desktop_file_path}")
            except Exception as e:
                QMessageBox.warning(
                    self, "Autostart Error", f"Failed to enable autostart on Linux: {e}"
                )
                logger.info(f"Linux autostart error: {e}")
        else:
            QMessageBox.warning(
                self,
                "Autostart",
                f"Autostart not supported on this operating system ({sys.platform}).",
            )

    def disable_autostart(self):
        if sys.platform.startswith("win"):
            # Windows: Remove the registry entry or delete the shortcut.
            # QMessageBox.information(self, "Autostart",
            #                         "On Windows, you'd typically remove the registry entry "
            #                         "or delete the shortcut from the Startup folder.")
            logger.info("Windows disable autostart.")
            import winreg

            try:
                key = winreg.OpenKey(
                    winreg.HKEY_CURRENT_USER,
                    r"Software\Microsoft\Windows\CurrentVersion\Run",
                    0,
                    winreg.KEY_SET_VALUE,
                )
                winreg.DeleteValue(key, "MyProgram")
                winreg.CloseKey(key)
                logger.info("Autostart disabled in Windows Registry.")
            except FileNotFoundError:
                logger.info(
                    "Autostart entry not found in Windows Registry (already disabled?)."
                )
            except Exception as e:
                QMessageBox.warning(
                    self,
                    "Autostart Error",
                    f"Failed to disable autostart on Windows: {e}",
                )
                logger.info(f"Windows disable autostart error: {e}")

        elif sys.platform.startswith("linux"):
            # Linux: Delete the .desktop file.
            desktop_file_path = os.path.expanduser(
                "~/.config/autostart/myprogram.desktop"
            )
            try:
                if os.path.exists(desktop_file_path):
                    os.remove(desktop_file_path)
                    logger.info(
                        f"Autostart .desktop file removed from: {desktop_file_path}"
                    )
                else:
                    logger.info(
                        "Autostart .desktop file not found (already disabled?)."
                    )
            except Exception as e:
                QMessageBox.warning(
                    self,
                    "Autostart Error",
                    f"Failed to disable autostart on Linux: {e}",
                )
                logger.info(f"Linux disable autostart error: {e}")
        else:
            pass  # Already handled by enable_autostart message


class MainWindow(QMainWindow, QtStyleTools):
    """
    The main application window for the Hybrid Search GUI.
    """

    def __init__(self):
        super().__init__()
        self.settings = QApplication.instance().settings

        self.setWindowTitle("Local Content Search")
        geometry = self.settings.value("window_geometry", QByteArray())
        if not geometry.isEmpty():
            self.restoreGeometry(geometry)
        else:
            self.setGeometry(100, 100, 800, 600)  # Initial window size
        # Ensure the main window is resizable: set a reasonable minimum and no restrictive maximum
        self.setMinimumSize(600, 400)
        self.setMaximumSize(16777215, 16777215)
        # Allow central widget to expand so resizing the window changes height/width
        self.central_widget = QWidget()
        self.setCentralWidget(self.central_widget)
        self.central_widget.setSizePolicy(QSizePolicy.Expanding, QSizePolicy.Expanding)
        # recreate main layout after resetting central widget
        self.main_layout = QVBoxLayout(self.central_widget)
        icon = QIcon("./resources/icon.png")
        self.setWindowIcon(icon)

        self.tray_icon = QSystemTrayIcon(icon, self)
        self.tray_icon.setToolTip("Local Content Search")
        tray_menu = QMenu()
        restore_action = QAction("Restor Window", self)
        quit_action = QAction("Exit", self)
        restore_action.triggered.connect(self.show)
        quit_action.triggered.connect(self.close_app)
        tray_menu.addAction(restore_action)
        tray_menu.addAction(quit_action)
        self.tray_icon.setContextMenu(tray_menu)
        self.tray_icon.activated.connect(self.on_tray_icon_activated)
        self.tray_icon.show()

        self.hybrid_search = HybridSearch()
        self.hybrid_search.index_rebuild_progress_signal.connect(
            self.on_rebuild_signal_progress
        )
        self.index_rebuilder_thread = None
        self.file_monitor_thread = None
        # Store last search results so filters can be applied locally without re-searching
        self.last_search_results = None
        self.last_highlight_tokens = None
        self.last_query = ""

        self._load_config()  # Load paths and last search mode
        self.hybrid_search.file_paths = set(self.config.get("file_paths", []))
        self.hybrid_search.folder_paths = set(self.config.get("folder_paths", []))

        self._setup_ui()
        self._setup_threads()  # Initialize threads but don't start file_monitor_thread yet

        # Attempt to load index on startup
        if self.hybrid_search.load_index_components():
            self.status_bar.showMessage("Index loaded successfully.")
            # Start file monitor after successful load
            self.file_monitor_thread.start()
        else:
            # If load fails or no index exists, rebuild
            self._start_rebuild_index()

        # show the help window on first launch
        is_first_time_launch = self.settings.value("is_first_time_launch", 1)
        if is_first_time_launch:
            self.settings_tab.on_help_button_pushed()
            self.settings.setValue("is_first_time_launch", 0)

    def _load_config(self):
        """Loads configuration from a JSON file."""
        self.config = {}
        if os.path.exists(CONFIG_FILE):
            with open(CONFIG_FILE, "r", encoding="utf-8") as f:
                self.config = json.load(f)
        else:
            logger.info(f"Config file not found: {CONFIG_FILE}")

    def _save_config(self):
        """Saves current configuration to a JSON file."""
        self.config["file_paths"] = list(self.hybrid_search.file_paths)
        self.config["folder_paths"] = list(self.hybrid_search.folder_paths)
        last_search_mode = list(SEARCH_MODE_MAPPING.values())[
            self.search_mode_combo.currentIndex()
        ]
        self.config["last_search_mode"] = last_search_mode
        with open(CONFIG_FILE, "w", encoding="utf-8") as f:
            json.dump(self.config, f, indent=4, ensure_ascii=False)

    def _setup_ui(self):
        """Sets up the main user interface."""
        self.tab_widget = QTabWidget()
        self.main_layout.addWidget(self.tab_widget)

        self._setup_search_tab()
        self._setup_management_tab()
        self._setup_settings_tab()
        self.status_bar = QStatusBar()
        self.setStatusBar(self.status_bar)
        self.status_bar.showMessage("Ready")

    def _setup_search_tab(self):
        """Sets up the Search tab UI."""
        self.search_tab = QWidget()
        self.tab_widget.addTab(self.search_tab, "Search")
        search_layout = QVBoxLayout(self.search_tab)

        # Search Input and Mode
        search_input_layout = QHBoxLayout()
        self.query_input = QLineEdit()
        self.query_input.setClearButtonEnabled(True)
        self.query_input.setPlaceholderText("Enter your query...")
        self.query_input.returnPressed.connect(self._perform_search)
        search_input_layout.addWidget(self.query_input)

        self.search_mode_combo = QComboBox()
        self.search_mode_combo.addItems(list(SEARCH_MODE_MAPPING.keys()))
        # Set last used search mode from config
        last_mode = self.config.get("last_search_mode", "hybrid")
        index = list(SEARCH_MODE_MAPPING.values()).index(last_mode)
        self.search_mode_combo.setCurrentIndex(index)
        search_input_layout.addWidget(self.search_mode_combo)

        self.search_button = QPushButton("Search")
        self.search_button.clicked.connect(self._perform_search)
        search_input_layout.addWidget(self.search_button)

        search_layout.addLayout(search_input_layout)

        # Main content area with sidebar and results table
        content_layout = QHBoxLayout()

        # Store default sidebar width for animation
        self.sidebar_default_width = 200

        # Create toggle button early so it can be used in results area
        self.toggle_sidebar_button = QPushButton("◀")  # Left arrow to collapse
        self.toggle_sidebar_button.setMaximumWidth(30)
        self.toggle_sidebar_button.clicked.connect(self._toggle_sidebar)

        # Left Sidebar for Filters (wrapped in a scroll area so items can be scrolled when space is small)
        self.sidebar = QScrollArea()
        self.sidebar.setWidgetResizable(True)
        sidebar_content = QWidget()
        sidebar_layout = QVBoxLayout(sidebar_content)
        # make sidebar more compact
        sidebar_layout.setSpacing(1)
        sidebar_layout.setContentsMargins(1, 1, 1, 1)

        # Sidebar header with toggle button
        sidebar_header_layout = QHBoxLayout()
        sidebar_title = QLabel("Filters")
        sidebar_title.setStyleSheet("font-weight: bold; font-size: 12px;")
        sidebar_header_layout.addWidget(sidebar_title)
        sidebar_header_layout.addStretch()
        sidebar_layout.addLayout(sidebar_header_layout)
        sidebar_layout.addSpacing(5)

        # Date Filter Section
        date_label = QLabel("Filter by Date:")
        date_label.setStyleSheet("font-weight: bold;")
        sidebar_layout.addWidget(date_label)

        self.date_filters = {}
        date_options = [
            ("All Time", None),
            ("Past 24 Hours", timedelta(hours=24)),
            ("Past 7 Days", timedelta(days=7)),
            ("Past 30 Days", timedelta(days=30)),
            ("Past 90 Days", timedelta(days=90)),
        ]

        for label, delta in date_options:
            checkbox = QCheckBox(label)
            checkbox.stateChanged.connect(self._on_filter_changed)
            self.date_filters[label] = (checkbox, delta)
            sidebar_layout.addWidget(checkbox)

        sidebar_layout.addSpacing(1)

        # File Type Filter Section
        file_type_label = QLabel("Filter by File Type:")
        file_type_label.setStyleSheet("font-weight: bold;")
        sidebar_layout.addWidget(file_type_label)

        self.file_type_filters = {}
        file_types = [
            ("All Types", None),
            ("PowerPoint", [".pptx"]),
            ("Excel", [".xlsx", ".xls"]),
            ("Word", [".docx"]),
            ("PDF", [".pdf"]),
            ("Text", [".txt"]),
            ("Markdown", [".md"]),
            ("CSV", [".csv"]),
            ("HTML", [".html", ".htm"]),
            ("ODT", [".odt"]),
            ("XML", [".xml"]),
        ]

        for label, extensions in file_types:
            checkbox = QCheckBox(label)
            checkbox.stateChanged.connect(self._on_filter_changed)
            self.file_type_filters[label] = (checkbox, extensions)
            sidebar_layout.addWidget(checkbox)

        # Set "All Time" as default
        self.date_filters["All Time"][0].setChecked(True)

        # Set "All Types" as default
        self.file_type_filters["All Types"][0].setChecked(True)

        sidebar_layout.addSpacing(1)

        # Clear Filters Button
        self.clear_filters_button = QPushButton("Clear Filters")
        self.clear_filters_button.clicked.connect(self._clear_filters)
        sidebar_layout.addWidget(self.clear_filters_button)

        sidebar_layout.addStretch()
        sidebar_content.setLayout(sidebar_layout)
        # Set sidebar width on the scroll area
        self.sidebar.setWidget(sidebar_content)
        self.sidebar.setMaximumWidth(240)

        # Restore sidebar visibility state from settings
        sidebar_visible = self.settings.value("sidebar_visible", True, type=bool)
        if not sidebar_visible:
            self.sidebar.setMaximumWidth(0)
            self.toggle_sidebar_button.setText("▶")

        # Create a container for sidebar with overlapping toggle button
        sidebar_container = QWidget()
        sidebar_container_layout = QVBoxLayout(sidebar_container)
        sidebar_container_layout.setContentsMargins(0, 0, 0, 0)
        sidebar_container_layout.setSpacing(0)
        sidebar_container_layout.addWidget(self.sidebar)

        # Add toggle button with negative margin to overlap sidebar edge
        sidebar_container_layout.addWidget(self.toggle_sidebar_button)

        content_layout.addWidget(sidebar_container)

        # Results area
        # results_area_layout = QHBoxLayout()

        # Search Results Table
        self.results_table = QTableWidget()
        table_headers = [
            "Filename",
            "Folder Path",
            "Snippet",
            "File Size",
            "Last Modified",
        ]
        self.results_table.setColumnCount(len(table_headers))
        self.results_table.setHorizontalHeaderLabels(table_headers)
        self.results_table.horizontalHeader().setSectionResizeMode(
            QHeaderView.ResizeMode.Interactive
        )
        self.results_table.horizontalHeader().setSectionsMovable(True)
        self.results_table.setWordWrap(True)
        header_state = self.settings.value("headerState")
        if header_state:
            self.results_table.horizontalHeader().restoreState(header_state)
        else:
            self.results_table.setColumnWidth(0, 180)
            self.results_table.setColumnWidth(1, 180)
            self.results_table.setColumnWidth(2, 280)

        # Resize columns to fit contents
        # self.results_table.resizeColumnsToContents()
        # self.results_table.setSortingEnabled(True) # Causes incomplete display
        self.results_table.setSortingEnabled(False)
        self.results_table.setSelectionBehavior(
            QAbstractItemView.SelectionBehavior.SelectItems
        )
        # self.results_table.setEditTriggers(QAbstractItemView.EditTrigger.NoEditTriggers)  # Make cells read-only
        self.results_table.setEditTriggers(
            QAbstractItemView.EditTrigger.SelectedClicked
        )
        self.results_table.cellDoubleClicked.connect(self._handle_table_double_click)
        # Enable custom context menu for result rows
        self.results_table.setContextMenuPolicy(Qt.CustomContextMenu)
        self.results_table.customContextMenuRequested.connect(
            self._show_results_context_menu
        )
        # Connect selection to preview
        self.results_table.itemSelectionChanged.connect(
            self._on_result_selection_changed
        )

        # Create a splitter for results table and preview pane (left=results, right=preview)
        self.splitter = QSplitter(Qt.Horizontal)
        self.splitter.addWidget(self.results_table)

        # Preview pane
        preview_container = QWidget()
        preview_layout = QVBoxLayout(preview_container)
        preview_label = QLabel("Preview")
        preview_label.setStyleSheet("font-weight: bold; font-size: 11px; margin: 5px;")
        preview_layout.addWidget(preview_label)

        self.preview_pane = UniversalPreviewPane()
        self.preview_pane.setMinimumWidth(20)
        preview_layout.addWidget(self.preview_pane)
        preview_layout.setContentsMargins(0, 0, 0, 0)
        preview_layout.setSpacing(0)

        self.splitter.addWidget(preview_container)
        self.splitter.setStretchFactor(0, 1)
        self.splitter.setStretchFactor(1, 1)

        # Store splitter state
        splitter_state = self.settings.value("splitterState")
        if splitter_state:
            self.splitter.restoreState(splitter_state)

        content_layout.addWidget(self.splitter)
        content_layout.setStretch(1, 1)
        search_layout.addLayout(content_layout)

    def _setup_management_tab(self):
        """Sets up the Management tab UI."""
        self.management_tab = QWidget()
        self.tab_widget.addTab(self.management_tab, "Search and Index Management")
        management_layout = QVBoxLayout(self.management_tab)

        # Folder Paths
        folder_layout = QVBoxLayout()
        folder_layout.addWidget(QLabel("Indexed Folders"))
        self.folder_list = QListWidget()
        for path in self.hybrid_search.folder_paths:
            self.folder_list.addItem(path)
        folder_layout.addWidget(self.folder_list)

        folder_buttons_layout = QHBoxLayout()
        self.add_folder_button = QPushButton("Add Folder")
        self.add_folder_button.clicked.connect(self._add_folder)
        folder_buttons_layout.addWidget(self.add_folder_button)

        self.remove_folder_button = QPushButton("Remove Selected Folder")
        self.remove_folder_button.clicked.connect(self._remove_folder)
        folder_buttons_layout.addWidget(self.remove_folder_button)
        folder_layout.addLayout(folder_buttons_layout)
        management_layout.addLayout(folder_layout)

        # File Paths
        file_layout = QVBoxLayout()
        file_layout.addWidget(QLabel("Indexed Files"))
        self.file_list = QListWidget()
        for path in self.hybrid_search.file_paths:
            self.file_list.addItem(path)
        file_layout.addWidget(self.file_list)

        file_buttons_layout = QHBoxLayout()
        self.add_file_button = QPushButton("Add File")
        self.add_file_button.clicked.connect(self._add_file)
        file_buttons_layout.addWidget(self.add_file_button)

        self.remove_file_button = QPushButton("Remove Selected File")
        self.remove_file_button.clicked.connect(self._remove_file)
        file_buttons_layout.addWidget(self.remove_file_button)
        file_layout.addLayout(file_buttons_layout)
        management_layout.addLayout(file_layout)

        # Rebuild Index Button
        self.rebuild_index_button = QPushButton("Rebuild Index Now")
        self.rebuild_index_button.clicked.connect(self._start_rebuild_index)
        management_layout.addWidget(self.rebuild_index_button)

        management_layout.addStretch()  # Push content to top

    def _setup_settings_tab(self):
        self.settings_tab = SettingsTab()
        self.tab_widget.addTab(self.settings_tab, "Settings")

    def _setup_threads(self):
        """Sets up the background threads for indexing and file monitoring."""
        # Index Rebuilder Thread
        self.index_rebuilder_thread = IndexRebuilderThread(self.hybrid_search)
        self.index_rebuilder_thread.rebuild_started.connect(self._on_rebuild_started)
        self.index_rebuilder_thread.rebuild_finished.connect(self._on_rebuild_finished)

        # File Monitor Thread (will be started after initial index load/rebuild)
        interval = self.settings.value("monitor_interval", 300, type=int)
        self.file_monitor_thread = FileMonitorThread(
            self.hybrid_search, interval=interval
        )  # Check every 30 seconds
        self.file_monitor_thread.changes_detected.connect(
            self._on_file_changes_detected
        )
        self.file_monitor_thread.monitoring_status_update.connect(
            self.status_bar.showMessage
        )

    @Slot()
    def _handle_table_double_click(self, row, column):
        item_widget = self.results_table.item(row, 1)
        item = item_widget.text() if item_widget is not None else None
        if not item:
            return
        if column == 1:  # open the folder
            if os.path.exists(item):
                QDesktopServices.openUrl(QUrl.fromLocalFile(item))
        else:  # open the file
            if os.path.exists(item):
                QDesktopServices.openUrl(QUrl.fromLocalFile(item))

    @Slot("QPoint")
    def _show_results_context_menu(self, pos):
        """Show context menu for a result row with actions: Open folder, Copy path."""
        row = self.results_table.rowAt(pos.y())
        if row < 0:
            return

        item_widget = self.results_table.item(row, 1)
        path = item_widget.text() if item_widget is not None else None
        if not path:
            return

        menu = QMenu(self)
        open_folder_act = QAction("Open folder", self)
        copy_path_act = QAction("Copy path", self)
        open_folder_act.triggered.connect(lambda: self._open_folder_for_row(path))
        copy_path_act.triggered.connect(lambda: self._copy_path_for_row(path))
        menu.addAction(open_folder_act)
        menu.addAction(copy_path_act)
        menu.exec(self.results_table.viewport().mapToGlobal(pos))

    def _open_folder_for_row(self, path):
        folder = str(Path(path))
        if os.path.exists(folder):
            QDesktopServices.openUrl(QUrl.fromLocalFile(folder))
        else:
            QMessageBox.warning(self, "Open folder", f"Path not found: {folder}")

    def _copy_path_for_row(self, path):
        try:
            QApplication.instance().clipboard().setText(path)
            self.status_bar.showMessage("Path copied to clipboard", 2000)
        except Exception:
            QMessageBox.information(self, "Copy path", path)

    @Slot()
    def _toggle_sidebar(self):
        """Toggle sidebar visibility with animation."""
        target_width = 0 if self.sidebar.width() > 50 else self.sidebar_default_width

        # Create animation for sidebar width
        self.sidebar_animation = QPropertyAnimation(self.sidebar, b"maximumWidth")
        self.sidebar_animation.setDuration(300)  # 300ms animation
        self.sidebar_animation.setStartValue(self.sidebar.maximumWidth())
        self.sidebar_animation.setEndValue(target_width)
        self.sidebar_animation.start()

        # Update button text and save state
        if target_width == 0:
            self.toggle_sidebar_button.setText("▶")  # Right arrow to expand
            self.settings.setValue("sidebar_visible", False)
        else:
            self.toggle_sidebar_button.setText("◀")  # Left arrow to collapse
            self.settings.setValue("sidebar_visible", True)

    @Slot()
    def _on_result_selection_changed(self):
        """Load preview when a result is selected."""
        selected_items = self.results_table.selectedIndexes()
        if not selected_items:
            self.preview_pane.clear()
            return

        # Get the file path from the selected row (column 1 is File Path)
        row = selected_items[0].row()
        file_name_item = self.results_table.item(row, 0)
        folder_path_item = self.results_table.item(row, 1)
        if not file_name_item or not folder_path_item:
            return

        file_path = Path(folder_path_item.text()) / file_name_item.text()
        self.preview_pane.handle_preview(str(file_path))

        """Load and display file preview."""
        try:
            if not os.path.exists(file_path):
                self.preview_pane.setText(f"File not found: {file_path}")
                return

            file_size = os.path.getsize(file_path)
            # Limit preview to several MB to avoid loading huge files
            FILE_SIZE_LIMIT_MB = 5
            if file_size > FILE_SIZE_LIMIT_MB * 1024 * 1024:
                self.preview_pane.setText(f"File is too large to preview)")
                return

            _, ext = os.path.splitext(file_path)
            ext = ext.lower()

            # Text-based files
            if ext in [
                ".txt",
                ".md",
                ".py",
                ".js",
                ".json",
                ".xml",
                ".html",
                ".htm",
                ".css",
                ".csv",
                ".log",
                ".sh",
                ".bat",
                ".yaml",
                ".yml",
                ".ini",
                ".conf",
                ".config",
            ]:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()
                    self.preview_pane.setText(content)

            # Office files - try to extract text
            elif ext in [".docx", ".doc"]:
                try:
                    from docx import Document

                    doc = Document(file_path)
                    text = "\n".join([para.text for para in doc.paragraphs])
                    self.preview_pane.setText(
                        text if text else "(No text content found)"
                    )
                except Exception:
                    self.preview_pane.setText(
                        "Preview not available for Word documents (python-docx not installed)"
                    )

            elif ext in [".xlsx", ".xls"]:
                try:
                    import openpyxl

                    wb = openpyxl.load_workbook(file_path, read_only=True)
                    sheet = wb.active
                    lines = []
                    for row in sheet.iter_rows(values_only=True):
                        lines.append(
                            "\t".join(
                                str(cell) if cell is not None else "" for cell in row
                            )
                        )
                    self.preview_pane.setText(
                        "\n".join(lines[:100]) if lines else "(No data)"
                    )
                except Exception:
                    self.preview_pane.setText(
                        "Preview not available for Excel files (openpyxl not installed)"
                    )

            elif ext in [".pptx"]:
                try:
                    from pptx import Presentation

                    prs = Presentation(file_path)
                    text_content = []
                    for slide_idx, slide in enumerate(prs.slides):
                        text_content.append(f"--- Slide {slide_idx + 1} ---")
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_content.append(shape.text)
                    self.preview_pane.setText(
                        "\n".join(text_content)
                        if text_content
                        else "(No text content found)"
                    )
                except Exception:
                    self.preview_pane.setText(
                        "Preview not available for PowerPoint files (python-pptx not installed)"
                    )

            # PDF files
            elif ext == ".pdf":
                try:
                    import PyPDF2

                    with open(file_path, "rb") as f:
                        pdf_reader = PyPDF2.PdfReader(f)
                        text_content = []
                        for page_idx, page in enumerate(
                            pdf_reader.pages[:5]
                        ):  # First 5 pages
                            text = page.extract_text()
                            if text:
                                text_content.append(
                                    f"--- Page {page_idx + 1} ---\n{text}"
                                )
                        self.preview_pane.setText(
                            "\n".join(text_content)
                            if text_content
                            else "(No text content found)"
                        )
                except Exception:
                    self.preview_pane.setText(
                        "Preview not available for PDF files (PyPDF2 not installed)"
                    )

            # Image files
            elif ext in [".jpg", ".jpeg", ".png", ".bmp", ".gif", ".tiff"]:
                self.preview_pane.setText(
                    f"Image file: {file_path}\n(Image preview not yet supported in text pane)"
                )

            else:
                self.preview_pane.setText(f"Preview not supported for file type: {ext}")

        except Exception as e:
            self.preview_pane.setText(f"Error loading preview: {str(e)}")

    @Slot()
    def _add_folder(self):
        """Opens a dialog to select a folder and adds it to the list."""
        folder_path = QFileDialog.getExistingDirectory(self, "Select Folder to Index")
        if folder_path and folder_path not in self.hybrid_search.folder_paths:
            self.hybrid_search.folder_paths.add(folder_path)
            self.folder_list.addItem(folder_path)
            self._save_config()  # Save config will also save the index

    @Slot()
    def _remove_folder(self):
        """Removes the selected folder from the list."""
        selected_items = self.folder_list.selectedItems()
        if not selected_items:
            return
        reply = QMessageBox.question(
            self,
            "Remove folder from monitoring",
            "Confirm removing the selected folder(s)? This will trigger index rebuilding.",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
        )
        if reply == QMessageBox.StandardButton.Yes:
            for item in selected_items:
                self.hybrid_search.folder_paths.discard(item.text())
                self.folder_list.takeItem(self.folder_list.row(item))
            self._save_config()  # Save config will also save the index

    @Slot()
    def _add_file(self):
        """Opens a dialog to select a file and adds it to the list."""
        file_path, _ = QFileDialog.getOpenFileName(self, "Select File to Index", "")
        if file_path and file_path not in self.hybrid_search.file_paths:
            self.hybrid_search.file_paths.add(file_path)
            self.file_list.addItem(file_path)
            self._save_config()  # Save config will also save the index

    @Slot()
    def _remove_file(self):
        """Removes the selected file from the list."""
        selected_items = self.file_list.selectedItems()
        if not selected_items:
            return
        reply = QMessageBox.question(
            self,
            "Remove file from monitoring",
            "Confirm removing the selected file(s)? This will trigger index rebuilding.",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No,
        )
        if reply == QMessageBox.StandardButton.Yes:
            for item in selected_items:
                self.hybrid_search.file_paths.discard(item.text())
                self.file_list.takeItem(self.file_list.row(item))
            self._save_config()  # Save config will also save the index

    @Slot()
    def _start_rebuild_index(self):
        """Starts the index rebuilding process in a separate thread."""
        if self.index_rebuilder_thread.isRunning():
            QMessageBox.information(
                self,
                "Indexing in Progress",
                "Index rebuilding is already in progress. Please wait.",
            )
            return

        self.status_bar.showMessage("Starting to rebuild index...")
        self.rebuild_index_button.setEnabled(False)
        self.search_button.setEnabled(False)
        self.index_rebuilder_thread.start()

    @Slot()
    def _on_rebuild_started(self):
        """Updates UI when index rebuilding starts."""
        self.status_bar.showMessage("Index is being rebuilt... Might take a while.")
        logger.info("Rebuild started signal received.")
        self.rebuild_start_time = time.time()

    @Slot(bool)
    def _on_rebuild_finished(self, success):
        """Updates UI when index rebuilding finishes."""
        if success:
            time_elapsed = time.time() - self.rebuild_start_time
            self.status_bar.showMessage(
                f"Index rebuild finished in {time_elapsed:.1f} seconds"
            )
            logger.info(
                f"Rebuild finished signal received: Success. Time used {time_elapsed} seconds"
            )

        else:
            self.status_bar.showMessage("Index rebuild failed.")
            logger.info("Rebuild finished signal received: Failure.")
        self.rebuild_index_button.setEnabled(True)
        self.search_button.setEnabled(True)

        # Ensure file monitor is always running and up-to-date after an index operation (load or rebuild)
        if self.file_monitor_thread and self.file_monitor_thread.isRunning():
            self.file_monitor_thread.stop()  # Stop if already running to restart with fresh mtimes
        # Re-initialize the thread to ensure it picks up latest paths and then start it
        interval = self.settings.value("monitor_interval", 300, type=int)
        self.file_monitor_thread = FileMonitorThread(
            self.hybrid_search, interval=interval
        )
        self.file_monitor_thread.changes_detected.connect(
            self._on_file_changes_detected
        )
        self.file_monitor_thread.monitoring_status_update.connect(
            self.status_bar.showMessage
        )
        self.file_monitor_thread.start()

    @Slot()
    def _on_file_changes_detected(self):
        logger.info("File changes detected! Triggering index rebuild.")
        self.status_bar.showMessage(
            "File changes detected! Triggering index rebuild..."
        )
        self._start_rebuild_index()

    @Slot()
    def _perform_search(self):
        """Performs a search based on the current query and selected mode."""
        query = self.query_input.text().strip()
        if not query:
            QMessageBox.warning(self, "Empty Query", "Please enter a search query.")
            return

        if self.index_rebuilder_thread.isRunning():
            QMessageBox.information(
                self,
                "Indexing in Progress",
                "Cannot search while index is being rebuilt. Please wait.",
            )
            return

        selected_mode = SEARCH_MODE_MAPPING[self.search_mode_combo.currentText()]
        logger.info(f"Searching for '{query}' with mode: {selected_mode}")
        self.status_bar.showMessage(f"Searching for '{query}'...")

        results, highlight_tokens = self.hybrid_search.search(query, mode=selected_mode)
        # store raw results so filters can be applied locally without triggering new searches
        self.last_search_results = results
        self.last_highlight_tokens = highlight_tokens
        self.last_query = query
        self._update_filter_counts(self.last_search_results)
        filtered_results = self._apply_filters(self.last_search_results)
        self._display_results(
            filtered_results, self.last_highlight_tokens, self.last_query
        )
        self.status_bar.showMessage(
            f"Search completed. Found {len(filtered_results)} results"
        )
        self._save_config()  # Save the last used search mode

    def _apply_filters(self, results):
        """Applies selected filters to search results."""
        filtered_results = results

        # Apply date filters
        selected_date_filter = None
        for label, (checkbox, delta) in self.date_filters.items():
            if checkbox.isChecked() and label != "All Time":
                selected_date_filter = delta
                break

        if selected_date_filter:
            cutoff_time = datetime.now() - selected_date_filter
            filtered_results = [
                result
                for result in filtered_results
                if result.get("mtime")
                and datetime.fromtimestamp(result["mtime"]) >= cutoff_time
            ]

        # Apply file type filters
        selected_file_types = []
        for label, (checkbox, extensions) in self.file_type_filters.items():
            if checkbox.isChecked() and label != "All Types":
                if extensions:
                    selected_file_types.extend(extensions)

        if selected_file_types:
            filtered_results = [
                result
                for result in filtered_results
                if any(
                    result.get("path", "").lower().endswith(ext)
                    for ext in selected_file_types
                )
            ]

        return filtered_results

    @Slot()
    def _on_filter_changed(self):
        """Keep "All" options exclusive with specifics, then refresh counts/results.

        Prioritize direct clicks on the "All" checkbox so it immediately takes effect
        (instead of being overridden by the presence of previously-checked specifics).
        """
        sender = self.sender()

        def enforce_exclusive(filters, all_key):
            all_cb = filters[all_key][0]
            # If the user clicked the "All" checkbox, honor that action first
            if sender is all_cb:
                if all_cb.isChecked():
                    # Uncheck all specific options
                    for k, (cb, _) in filters.items():
                        if k != all_key:
                            cb.blockSignals(True)
                            cb.setChecked(False)
                            cb.blockSignals(False)
                else:
                    # If user unchecked All and no specific is selected, keep All checked
                    any_specific = any(
                        cb.isChecked() for k, (cb, _) in filters.items() if k != all_key
                    )
                    if not any_specific:
                        all_cb.blockSignals(True)
                        all_cb.setChecked(True)
                        all_cb.blockSignals(False)
            else:
                # A specific checkbox changed: if any specific is checked -> uncheck All, else check All
                specific_checked = any(
                    cb.isChecked() for k, (cb, _) in filters.items() if k != all_key
                )
                all_cb.blockSignals(True)
                all_cb.setChecked(not specific_checked)
                all_cb.blockSignals(False)

        enforce_exclusive(self.date_filters, "All Time")
        enforce_exclusive(self.file_type_filters, "All Types")

        # Reapply filters using the last stored search results (no new search)
        if self.last_search_results is not None:
            self._update_filter_counts(self.last_search_results)
            filtered = self._apply_filters(self.last_search_results)
            self._display_results(filtered, self.last_highlight_tokens, self.last_query)

    @Slot()
    def _clear_filters(self):
        """Clears all filters and resets to default."""
        # Reset date filters to "All Time"
        for label, (checkbox, _) in self.date_filters.items():
            checkbox.blockSignals(True)
            checkbox.setChecked(label == "All Time")
            checkbox.blockSignals(False)

        # Reset file type filters to "All Types"
        for label, (checkbox, _) in self.file_type_filters.items():
            checkbox.blockSignals(True)
            checkbox.setChecked(label == "All Types")
            checkbox.blockSignals(False)

        # Reapply filters using stored last search results (do not trigger a new search)
        if self.last_search_results is not None:
            self._update_filter_counts(self.last_search_results)
            filtered = self._apply_filters(self.last_search_results)
            self._display_results(filtered, self.last_highlight_tokens, self.last_query)

    def _update_filter_counts(self, results):
        """Updates the counts displayed next to filter checkboxes based on search results."""
        # Update date filter counts
        for label, (checkbox, delta) in self.date_filters.items():
            if label == "All Time":
                count = len(results)
            else:
                cutoff_time = datetime.now() - delta
                count = len(
                    [
                        r
                        for r in results
                        if r.get("mtime")
                        and datetime.fromtimestamp(r["mtime"]) >= cutoff_time
                    ]
                )
            checkbox.setText(f"{label} ({count})")

        # Update file type filter counts
        for label, (checkbox, extensions) in self.file_type_filters.items():
            if label == "All Types":
                count = len(results)
            else:
                if extensions:
                    count = len(
                        [
                            r
                            for r in results
                            if any(
                                r.get("path", "").lower().endswith(ext)
                                for ext in extensions
                            )
                        ]
                    )
                else:
                    count = 0
            checkbox.setText(f"{label} ({count})")

    def _display_results(self, results, highlight_tokens=None, query=""):
        """Displays search results in the QTableWidget. query_tokens used for highlighting"""
        self.results_table.setRowCount(0)  # Clear existing results
        highlight_words = highlight_tokens

        for row_idx, result in enumerate(results):
            self.results_table.insertRow(row_idx)
            file_name = os.path.basename(result.get("path", "N/A"))
            file_path = result.get("path", "N/A")
            content = result.get("content", "N/A")
            file_size = (
                f"{result.get('size', 0) / 1024:.2f} KB"
                if result.get("size") is not None
                else "N/A"
            )
            mtime_timestamp = result.get("mtime")
            date_modified = (
                datetime.fromtimestamp(mtime_timestamp).strftime("%Y-%m-%d %H:%M:%S")
                if mtime_timestamp
                else "N/A"
            )
            # source_mode = result.get('source_mode', 'N/A')

            _item = QTableWidgetItem(file_name)
            _item.setToolTip(file_name)
            self.results_table.setItem(row_idx, 0, _item)
            folder_path = os.path.dirname(file_path)
            # making the forward and backward slashes consistent
            folder_path = str(Path(folder_path))
            _item = QTableWidgetItem(folder_path)
            _item.setToolTip(folder_path)
            self.results_table.setItem(row_idx, 1, _item)
            # self.results_table.setCellWidget(row_idx, 2, label)
            # item.setTextAlignment(Qt.AlignLeft | Qt.AlignTop)

            # text_edit1.setReadOnly(True)
            text_edit = QTextEdit()
            text_edit.setHtml(create_highlighted_text(content, highlight_words, query))
            text_edit.setMaximumHeight(80)
            # text_edit.setStyleSheet('')
            # text_edit.setStyle(QApplication.style())
            self.results_table.setCellWidget(row_idx, 2, text_edit)

            # self.results_table.setItem(row_idx, 2, QTableWidgetItem(snippet))
            self.results_table.setItem(row_idx, 3, QTableWidgetItem(file_size))
            self.results_table.setItem(row_idx, 4, QTableWidgetItem(date_modified))
            # self.results_table.setItem(row_idx, 5, QTableWidgetItem(source_mode))

        self.results_table.setWordWrap(True)
        self.results_table.resizeRowsToContents()  # Adjust row heights to fit content

        # self.results_table.resizeColumnsToContents()  # Adjust column widths to fit content

    def on_tray_icon_activated(self, reason):
        if reason == QSystemTrayIcon.Trigger:  # Left click
            self.show()
            self.raise_()
            self.activateWindow()

    @Slot(str)
    def on_rebuild_signal_progress(self, file_name):
        self.status_bar.showMessage(f"Rebuilding index...Processing: {file_name}")

    def closeEvent(self, event):
        """Handles the window close event."""
        if self.settings.value("minimize_to_tray", True, type=bool):
            event.ignore()
            self.hide()
            # self.tray_icon.showMessage("Tray Application", "Application minimized to tray.",
            #                            QSystemTrayIcon.Information, 2000)
        else:
            self.tray_icon.hide()
            logger.info("Application closing. Stopping threads...")
            if self.file_monitor_thread and self.file_monitor_thread.isRunning():
                logger.info("Stopping file monitor thread")
                self.file_monitor_thread.stop()
            if self.index_rebuilder_thread and self.index_rebuilder_thread.isRunning():
                # self.index_rebuilder_thread.wait()
                logger.info("Stopping rebuilder thread")
                self.index_rebuilder_thread.stop()
            logger.info("Saving config")
            self._save_config()
            self.settings.setValue("window_geometry", self.saveGeometry())
            header = self.results_table.horizontalHeader()
            self.settings.setValue("headerState", header.saveState())
            self.settings.setValue("splitterState", self.splitter.saveState())
            self.settings.setValue("font_size", int(self.settings_tab.spin_box.value()))
            logger.info("Config saved")
            super().closeEvent(event)

    def close_app(self):
        self.tray_icon.hide()
        QApplication.instance().quit()


if __name__ == "__main__":
    app = QApplication(sys.argv)
    is_autostart = AUTOSTART_ARG in sys.argv
    logger.info(f"Application launched. Is autostart: {is_autostart}")
    app.settings = QSettings("LocalContentSearch", "LocalContentSearch")  # Organization and application name
    app.setWindowIcon(QIcon("./resources/icon.png"))
    font_id = QFontDatabase.addApplicationFont("./resources/SourceHanSansSC-Normal.otf")
    font_families = QFontDatabase.applicationFontFamilies(font_id)
    custom_font_family = font_families[0]
    font_size = app.settings.value("font_size", 12, type=int)
    custom_font = QFont(custom_font_family, font_size)
    app.setFont(custom_font)

    if is_autostart:
        pass
    else:
        # splash screen
        splash_pix = QPixmap(400, 300)
        splash_pix.fill(QColor("#ecf0f1"))  # Fill with white background
        splash = QSplashScreen(splash_pix)
        splash.showMessage(
            "Program initializing...\n \nStarting inference engine\n\n Loading models...",
            Qt.AlignCenter,
            Qt.black,
        )
        splash.show()
        app.processEvents()
    import_and_initialize()

    window = MainWindow()
    THEME_EXTRA_SETTINGS = {
        "density_scale": "-2",  # Lower density (more compact layout)
        "font_size": font_size,
        "font_family": custom_font_family,
    }
    apply_stylesheet(
        app,
        theme="light_blue.xml",
        invert_secondary=True,
        extra=THEME_EXTRA_SETTINGS,
        css_file="custom.css",
    )
    if is_autostart:
        window.setWindowState(Qt.WindowMinimized)
    else:
        window.show()
        window.raise_()
        window.activateWindow()
        splash.finish(window)
    sys.exit(app.exec())
